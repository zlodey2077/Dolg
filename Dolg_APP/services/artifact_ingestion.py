"""Engineering artifact ingestion for review, learning and local AI memory.

The service keeps heavy or optional parsers behind lazy imports. It returns a
stable normalized payload that can be stored in EngineeringArtifact.facts and
reused by ProjectReview, learning-by-artifact and neural training examples.
"""

from __future__ import annotations

import hashlib
import re
from collections import Counter
from pathlib import Path
from typing import Any

TEXT_LIMIT = 20000


EXTENSION_PARSERS = {
    '.docx': 'docx',
    '.pdf': 'pdf',
    '.pptx': 'pptx',
    '.dxf': 'dxf',
    '.net': 'pcad_net',
    '.drc': 'pcad_drc',
    '.erc': 'pcad_erc',
    '.doc': 'ole',
    '.vsd': 'ole',
    '.dwg': 'dwg_stub',
    '.ms14': 'multisim_stub',
    '.pcb': 'pcad_binary_stub',
    '.sch': 'pcad_binary_stub',
}


def checksum_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def normalize_artifact_payload(
    *,
    source_name: str,
    source_path: str = '',
    data: bytes,
    parser: str,
    artifact_type: str,
    status: str = 'parsed',
    summary: str = '',
    facts: dict[str, Any] | None = None,
    warnings: list[str] | None = None,
    errors: list[str] | None = None,
) -> dict[str, Any]:
    return {
        'ok': status != 'error',
        'source_name': source_name,
        'source_path': source_path,
        'artifact_type': artifact_type,
        'parser': parser,
        'status': status,
        'checksum': checksum_bytes(data),
        'size_bytes': len(data),
        'summary': summary[:2000],
        'facts': facts or {},
        'warnings': warnings or [],
        'errors': errors or [],
    }


def parse_artifact(path: str | Path) -> dict[str, Any]:
    path = Path(path)
    data = path.read_bytes()
    parser = EXTENSION_PARSERS.get(path.suffix.lower(), 'unknown')
    try:
        if parser == 'docx':
            return parse_docx(path, data)
        if parser == 'pdf':
            return parse_pdf(path, data)
        if parser == 'pptx':
            return parse_pptx(path, data)
        if parser == 'dxf':
            return parse_dxf(path, data)
        if parser == 'pcad_net':
            return parse_pcad_net(path, data)
        if parser == 'pcad_drc':
            return parse_pcad_drc(path, data)
        if parser == 'pcad_erc':
            return parse_pcad_erc(path, data)
        if parser == 'ole':
            return parse_ole_metadata(path, data)
        if parser == 'dwg_stub':
            return parse_dwg_stub(path, data)
        if parser == 'multisim_stub':
            return parse_multisim_stub(path, data)
        if parser == 'pcad_binary_stub':
            return parse_pcad_binary_stub(path, data)
        return normalize_artifact_payload(
            source_name=path.name,
            source_path=str(path),
            data=data,
            parser='unknown',
            artifact_type='unknown',
            status='unsupported',
            summary='Unsupported artifact format.',
            facts={'document_summary': {}, 'cad_artifact': {}, 'check_report': {}},
            warnings=['Unsupported artifact format; metadata only.'],
        )
    except Exception as exc:
        return normalize_artifact_payload(
            source_name=path.name,
            source_path=str(path),
            data=data,
            parser=parser,
            artifact_type='unknown',
            status='error',
            summary=f'Artifact parser failed: {exc}',
            facts={},
            errors=[str(exc)],
        )


def save_artifact_report(report: dict[str, Any], *, project=None, user=None):
    from Dolg_APP.models import EngineeringArtifact, ProjectEvent

    artifact, _created = EngineeringArtifact.objects.update_or_create(
        project=project,
        checksum=report.get('checksum', ''),
        source_name=report.get('source_name') or 'artifact',
        defaults={
            'user': user if getattr(user, 'is_authenticated', False) else None,
            'source_path': report.get('source_path', ''),
            'artifact_type': report.get('artifact_type', 'unknown'),
            'parser': report.get('parser', ''),
            'status': report.get('status', 'parsed'),
            'size_bytes': int(report.get('size_bytes') or 0),
            'summary': report.get('summary', ''),
            'facts': report.get('facts') or {},
            'warnings': report.get('warnings') or [],
            'errors': report.get('errors') or [],
        },
    )
    if project is not None:
        ProjectEvent.log(
            project=project,
            user=user,
            event_type='artifact_ingested',
            payload={
                'artifact_id': artifact.id,
                'source_name': artifact.source_name,
                'artifact_type': artifact.artifact_type,
                'status': artifact.status,
            },
        )
    return artifact


def parse_docx(path: Path, data: bytes) -> dict[str, Any]:
    from docx import Document

    document = Document(str(path))
    paragraphs = [_clean_text(item.text) for item in document.paragraphs if _clean_text(item.text)]
    headings = []
    for paragraph in document.paragraphs:
        style = getattr(getattr(paragraph, 'style', None), 'name', '') or ''
        text = _clean_text(paragraph.text)
        if text and style.lower().startswith('heading'):
            headings.append(text)
    text = '\n'.join(paragraphs)
    facts = {
        'document_summary': {
            'format': 'docx',
            'paragraph_count': len(paragraphs),
            'heading_count': len(headings),
            'headings': headings[:24],
            'keywords': _keywords(text),
            'text_excerpt': text[:TEXT_LIMIT],
        },
        'process_steps': _process_steps(text),
    }
    return normalize_artifact_payload(
        source_name=path.name,
        source_path=str(path),
        data=data,
        parser='docx',
        artifact_type='document',
        summary=_summary_from_text(text, fallback=f'DOCX document: {len(paragraphs)} paragraphs.'),
        facts=facts,
    )


def parse_pdf(path: Path, data: bytes) -> dict[str, Any]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = len(reader.pages)
    chunks = []
    for page in reader.pages[:30]:
        chunks.append(page.extract_text() or '')
    text = _clean_text('\n'.join(chunks))
    facts = {
        'document_summary': {
            'format': 'pdf',
            'page_count': pages,
            'keywords': _keywords(text),
            'text_excerpt': text[:TEXT_LIMIT],
        },
        'process_steps': _process_steps(text),
    }
    return normalize_artifact_payload(
        source_name=path.name,
        source_path=str(path),
        data=data,
        parser='pdf',
        artifact_type='document',
        summary=_summary_from_text(text, fallback=f'PDF document: {pages} pages.'),
        facts=facts,
    )


def parse_pptx(path: Path, data: bytes) -> dict[str, Any]:
    from pptx import Presentation

    deck = Presentation(str(path))
    slide_texts = []
    for slide in deck.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and _clean_text(shape.text):
                texts.append(_clean_text(shape.text))
        slide_texts.append('\n'.join(texts))
    text = '\n'.join(slide_texts)
    facts = {
        'document_summary': {
            'format': 'pptx',
            'slide_count': len(deck.slides),
            'keywords': _keywords(text),
            'slides': slide_texts[:30],
            'text_excerpt': text[:TEXT_LIMIT],
        },
        'process_steps': _process_steps(text),
    }
    return normalize_artifact_payload(
        source_name=path.name,
        source_path=str(path),
        data=data,
        parser='pptx',
        artifact_type='document',
        summary=_summary_from_text(text, fallback=f'PPTX deck: {len(deck.slides)} slides.'),
        facts=facts,
    )


def parse_dxf(path: Path, data: bytes) -> dict[str, Any]:
    import ezdxf

    doc = ezdxf.readfile(str(path))
    modelspace = doc.modelspace()
    entity_types = Counter()
    layers = Counter()
    text_items = []
    for entity in modelspace:
        entity_types[entity.dxftype()] += 1
        layers[getattr(entity.dxf, 'layer', '') or '0'] += 1
        if entity.dxftype() == 'TEXT':
            text_items.append(_clean_text(entity.dxf.text))
        elif entity.dxftype() == 'MTEXT':
            text_items.append(_clean_text(entity.text))
    text_items = [item for item in text_items if item]
    facts = {
        'cad_artifact': {
            'format': 'dxf',
            'dxf_version': doc.dxfversion,
            'entity_count': sum(entity_types.values()),
            'entity_types': dict(entity_types.most_common(20)),
            'layers': dict(layers.most_common(30)),
            'texts': text_items[:50],
        },
        'process_steps': _process_steps(' '.join(text_items)),
    }
    return normalize_artifact_payload(
        source_name=path.name,
        source_path=str(path),
        data=data,
        parser='dxf',
        artifact_type='cad_drawing',
        summary=f'DXF {doc.dxfversion}: {sum(entity_types.values())} entities, {len(layers)} layers.',
        facts=facts,
    )


def parse_pcad_net(path: Path, data: bytes) -> dict[str, Any]:
    text = _decode_bytes(data)
    component_blocks = re.findall(r'\[\s*(.*?)\s*\]', text, flags=re.S)
    net_blocks = re.findall(r'\(\s*(.*?)\s*\)', text, flags=re.S)
    components = []
    for block in component_blocks:
        lines = [line.strip() for line in block.splitlines() if line.strip()]
        if not lines:
            continue
        components.append(
            {
                'ref': lines[0],
                'package': lines[1] if len(lines) > 1 else '',
                'part_number': lines[2] if len(lines) > 2 else '',
            }
        )
    nets = []
    for block in net_blocks:
        lines = [line.strip() for line in block.splitlines() if line.strip()]
        if not lines:
            continue
        nets.append({'name': lines[0], 'pins': lines[1:]})
    facts = {
        'cad_artifact': {
            'format': 'pcad_net',
            'component_count': len(components),
            'net_count': len(nets),
            'components': components[:200],
            'nets': nets[:200],
        },
        'bom_artifact': {
            'items': components[:200],
            'item_count': len(components),
        },
    }
    return normalize_artifact_payload(
        source_name=path.name,
        source_path=str(path),
        data=data,
        parser='pcad_net',
        artifact_type='cad_netlist',
        summary=f'P-CAD netlist: {len(components)} components, {len(nets)} nets.',
        facts=facts,
    )


def parse_pcad_drc(path: Path, data: bytes) -> dict[str, Any]:
    text = _decode_bytes(data)
    findings = []
    for match in re.finditer(
        r'Error\s+(\d+)\s+--\s+(.*?)(?=\nError\s+\d+\s+--|\n[-=]{6,}|\Z)', text, flags=re.S
    ):
        number = int(match.group(1))
        raw = _clean_text(match.group(2))
        findings.append(_pcad_drc_finding(number, raw))
    warnings_count = _first_int(r'(\d+)\s+warning\(s\)\s+detected', text)
    errors_count = _first_int(r'(\d+)\s+error\(s\)\s+detected', text)
    categories = Counter(item['category'] for item in findings)
    facts = {
        'check_report': {
            'format': 'pcad_drc',
            'tool': 'P-CAD Design Rule Check',
            'warning_count': warnings_count,
            'error_count': errors_count if errors_count is not None else len(findings),
            'finding_count': len(findings),
            'categories': dict(categories),
            'findings': findings,
        },
        'fault_cases': [_fault_case_from_finding(item) for item in findings[:30]],
    }
    return normalize_artifact_payload(
        source_name=path.name,
        source_path=str(path),
        data=data,
        parser='pcad_drc',
        artifact_type='check_report',
        status='partial' if findings and errors_count not in (None, len(findings)) else 'parsed',
        summary=f'P-CAD DRC: {len(findings)} parsed findings.',
        facts=facts,
        warnings=[] if findings else ['No DRC findings parsed from report.'],
    )


def parse_pcad_erc(path: Path, data: bytes) -> dict[str, Any]:
    text = _decode_bytes(data)
    categories = {}
    current = ''
    for line in text.splitlines():
        line = line.strip()
        if line.endswith(':') and line.isupper():
            current = line.rstrip(':').lower().replace(' ', '_')
        if current and line.startswith('Errors:'):
            categories.setdefault(current, {})['errors'] = _first_int(r'Errors:\s*(\d+)', line) or 0
        if current and line.startswith('Warnings:'):
            categories.setdefault(current, {})['warnings'] = _first_int(r'Warnings:\s*(\d+)', line) or 0
    total_errors = sum(item.get('errors', 0) for item in categories.values())
    total_warnings = sum(item.get('warnings', 0) for item in categories.values())
    facts = {
        'check_report': {
            'format': 'pcad_erc',
            'tool': 'P-CAD Electrical Rules Check',
            'warning_count': total_warnings,
            'error_count': total_errors,
            'categories': categories,
            'findings': [],
        },
    }
    return normalize_artifact_payload(
        source_name=path.name,
        source_path=str(path),
        data=data,
        parser='pcad_erc',
        artifact_type='check_report',
        summary=f'P-CAD ERC: {total_errors} errors, {total_warnings} warnings.',
        facts=facts,
    )


def parse_ole_metadata(path: Path, data: bytes) -> dict[str, Any]:
    import olefile

    ole = olefile.OleFileIO(str(path))
    try:
        streams = ['/'.join(item) for item in ole.listdir(streams=True, storages=False)]
        metadata = ole.get_metadata()
        props = {}
        for key in (
            'title',
            'subject',
            'author',
            'last_saved_by',
            'creating_application',
            'num_pages',
            'num_words',
        ):
            value = getattr(metadata, key, None)
            if value:
                props[key] = _decode_metadata(value)
    finally:
        ole.close()
    artifact_type = 'cad_drawing' if path.suffix.lower() == '.vsd' else 'document'
    facts = {
        'document_summary': {
            'format': path.suffix.lower().lstrip('.'),
            'ole_streams': streams[:80],
            'metadata': props,
        },
    }
    return normalize_artifact_payload(
        source_name=path.name,
        source_path=str(path),
        data=data,
        parser='ole',
        artifact_type=artifact_type,
        status='partial',
        summary=f'OLE artifact: {len(streams)} streams; metadata keys: {len(props)}.',
        facts=facts,
        warnings=[
            'Only OLE metadata was extracted; content conversion is required for full text/diagram analysis.'
        ],
    )


def parse_dwg_stub(path: Path, data: bytes) -> dict[str, Any]:
    version = data[:6].decode('ascii', errors='ignore')
    facts = {
        'cad_artifact': {
            'format': 'dwg',
            'dwg_version': version,
            'requires_conversion': True,
            'recommended_conversion': 'DWG to DXF',
        },
    }
    return normalize_artifact_payload(
        source_name=path.name,
        source_path=str(path),
        data=data,
        parser='dwg_stub',
        artifact_type='cad_drawing',
        status='unsupported',
        summary=f'DWG {version}: metadata only; convert to DXF for entity extraction.',
        facts=facts,
        warnings=['DWG is a closed binary CAD format; convert to DXF before full ingestion.'],
    )


def parse_multisim_stub(path: Path, data: bytes) -> dict[str, Any]:
    signature = data[:40].decode('latin1', errors='ignore')
    facts = {
        'cad_artifact': {
            'format': 'ms14',
            'signature': signature,
            'requires_conversion': True,
            'recommended_conversion': 'Multisim export to netlist/SPICE or XML',
        },
    }
    return normalize_artifact_payload(
        source_name=path.name,
        source_path=str(path),
        data=data,
        parser='multisim_stub',
        artifact_type='simulation',
        status='unsupported',
        summary='Multisim compressed workbook: metadata only; export is required for full ingestion.',
        facts=facts,
        warnings=['MS14 is a compressed Multisim workbook; export to SPICE/netlist for full analysis.'],
    )


def parse_pcad_binary_stub(path: Path, data: bytes) -> dict[str, Any]:
    text = _decode_bytes(data)
    strings = []
    for item in re.findall(r'[\w$+\-().,;:/\\ ]{5,}', text):
        item = _clean_text(item)
        if item and item not in strings:
            strings.append(item)
        if len(strings) >= 80:
            break
    fmt = 'pcad_schematic_binary' if path.suffix.lower() == '.sch' else 'pcad_pcb_binary'
    facts = {
        'cad_artifact': {
            'format': fmt,
            'readable_strings': strings[:80],
            'requires_export': True,
            'recommended_export': 'P-CAD netlist, DRC/ERC report or DXF',
        },
    }
    return normalize_artifact_payload(
        source_name=path.name,
        source_path=str(path),
        data=data,
        parser='pcad_binary_stub',
        artifact_type='cad_drawing',
        status='partial',
        summary=f'P-CAD binary artifact: {len(strings)} readable strings extracted.',
        facts=facts,
        warnings=[
            'P-CAD binary content is partially readable; use netlist/DRC/ERC/DXF exports for full analysis.'
        ],
    )


def artifact_reports_from_project(project, *, limit: int = 12) -> list[dict[str, Any]]:
    if project is None or not getattr(project, 'pk', None):
        return []
    try:
        qs = project.engineering_artifacts.all()[:limit]
    except Exception:
        return []
    reports = []
    for artifact in qs:
        reports.append(
            {
                'source_name': artifact.source_name,
                'source_path': artifact.source_path,
                'artifact_type': artifact.artifact_type,
                'parser': artifact.parser,
                'status': artifact.status,
                'checksum': artifact.checksum,
                'size_bytes': artifact.size_bytes,
                'summary': artifact.summary,
                'facts': artifact.facts or {},
                'warnings': artifact.warnings or [],
                'errors': artifact.errors or [],
            }
        )
    return reports


def review_external_cad_artifacts(reports: list[dict[str, Any]] | None) -> dict[str, Any]:
    findings = []
    artifact_count = 0
    unsupported_count = 0
    for report in reports or []:
        artifact_count += 1
        if report.get('status') == 'unsupported':
            unsupported_count += 1
        facts = report.get('facts') or {}
        check_report = facts.get('check_report') or {}
        for finding in check_report.get('findings') or []:
            item = dict(finding)
            item['source_name'] = report.get('source_name')
            item.setdefault('severity', 'warning')
            findings.append(item)
    severity_rank = {'critical': 4, 'error': 3, 'risk': 2, 'warning': 1, 'info': 0}
    findings.sort(key=lambda item: severity_rank.get(item.get('severity'), 0), reverse=True)
    return {
        'artifact_count': artifact_count,
        'unsupported_count': unsupported_count,
        'finding_count': len(findings),
        'findings': findings[:100],
    }


def learning_tasks_from_artifact(report: dict[str, Any], *, limit: int = 4) -> list[dict[str, Any]]:
    tasks = []
    check_report = (report.get('facts') or {}).get('check_report') or {}
    for finding in check_report.get('findings') or []:
        category = finding.get('category') or 'drc'
        title = {
            'short': 'Find and explain a short circuit',
            'outside_board': 'Move elements back inside the board outline',
            'uncommitted_pin': 'Fix an uncommitted or unconnected pin',
        }.get(category, 'Explain an imported DRC finding')
        tasks.append(
            {
                'task_type': 'circuit_build',
                'title': title,
                'prompt': finding.get('title') or finding.get('raw') or title,
                'rubric': {
                    'source_artifact': report.get('source_name'),
                    'source_parser': report.get('parser'),
                    'source_rule_id': finding.get('rule_id'),
                    'expected_fix': finding.get('recommendation'),
                },
            }
        )
        if len(tasks) >= limit:
            break
    return tasks


def training_examples_from_artifact(report: dict[str, Any], *, limit: int = 10) -> list[dict[str, Any]]:
    examples = []
    facts = report.get('facts') or {}
    check_report = facts.get('check_report') or {}
    for finding in check_report.get('findings') or []:
        examples.append(
            {
                'kind': 'drc_finding',
                'prompt': f'Explain {finding.get("rule_id")} from {report.get("source_name")}',
                'target': finding.get('recommendation') or finding.get('title') or finding.get('raw') or '',
                'features': {
                    'source_name': report.get('source_name'),
                    'parser': report.get('parser'),
                    'finding': finding,
                },
            }
        )
        if len(examples) >= limit:
            return examples
    summary = report.get('summary') or ''
    if summary:
        examples.append(
            {
                'kind': 'artifact_summary',
                'prompt': f'Summarize engineering artifact {report.get("source_name")}',
                'target': summary,
                'features': {
                    'artifact_type': report.get('artifact_type'),
                    'parser': report.get('parser'),
                    'keywords': (facts.get('document_summary') or {}).get('keywords', []),
                },
            }
        )
    return examples[:limit]


def _pcad_drc_finding(number: int, raw: str) -> dict[str, Any]:
    lower = raw.lower()
    category = 'other'
    severity = 'warning'
    title = 'Imported DRC finding'
    recommendation = 'Inspect imported CAD evidence before release.'
    rule_id = 'external.pcad.drc'
    if 'shorted to net' in lower:
        category = 'short'
        severity = 'critical'
        title = 'Short circuit between nets'
        recommendation = 'Separate the shorted nets and rerun DRC.'
        rule_id = 'external.pcad.short'
    elif 'outside board outline' in lower:
        category = 'outside_board'
        severity = 'error'
        title = 'Element or trace outside board outline'
        recommendation = 'Move the element or route inside the board outline.'
        rule_id = 'external.pcad.outside_board'
    elif 'uncommitted pin' in lower:
        category = 'uncommitted_pin'
        severity = 'error'
        title = 'Net connected to uncommitted pin'
        recommendation = 'Check the symbol pin mapping and commit/connect the pin.'
        rule_id = 'external.pcad.uncommitted_pin'
    elif 'unrouted' in lower:
        category = 'unrouted'
        severity = 'warning'
        title = 'Unrouted net'
        recommendation = 'Finish routing or mark the net intentionally unrouted.'
        rule_id = 'external.pcad.unrouted'
    nets = re.findall(r'Net\s+([A-Za-z0-9_+\-]+)', raw)
    refs = sorted(set(re.findall(r'\b([A-ZА-Я]{1,3}\d+)-\d+\b', raw)))
    return {
        'rule_id': rule_id,
        'number': number,
        'category': category,
        'severity': severity,
        'title': title,
        'raw': raw[:1000],
        'evidence': {
            'nets': nets[:8],
            'references': refs[:12],
            'coordinates': re.findall(r'\(([-0-9.]+),([-0-9.]+)\)', raw)[:8],
        },
        'recommendation': recommendation,
    }


def _fault_case_from_finding(finding: dict[str, Any]) -> dict[str, Any]:
    return {
        'symptom': finding.get('title') or finding.get('category'),
        'cause': finding.get('raw', '')[:300],
        'check': 'Run imported DRC/ERC and inspect highlighted coordinates or references.',
        'action': finding.get('recommendation') or 'Inspect imported CAD evidence.',
        'source_rule_id': finding.get('rule_id'),
    }


def _decode_bytes(data: bytes) -> str:
    for encoding in ('utf-8', 'cp1251', 'latin1'):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode('latin1', errors='ignore')


def _decode_metadata(value: Any) -> Any:
    if isinstance(value, bytes):
        for encoding in ('utf-8', 'cp1251', 'latin1'):
            try:
                return value.decode(encoding)
            except UnicodeDecodeError:
                continue
        return value.decode('latin1', errors='ignore')
    return value


def _clean_text(value: str) -> str:
    return re.sub(r'\s+', ' ', str(value or '')).strip()


def _summary_from_text(text: str, *, fallback: str) -> str:
    text = _clean_text(text)
    if not text:
        return fallback
    sentences = re.split(r'(?<=[.!?])\s+', text)
    summary = ' '.join(sentences[:2]).strip()
    return summary[:600] or fallback


def _keywords(text: str) -> list[str]:
    lowered = str(text or '').lower()
    candidates = {
        'requirements': ('requirement', 'требован', 'sys.', 'prog.', 'trace'),
        'schematic': ('schematic', 'схем', 'netlist', 'erc', 'drc'),
        'simulation': ('simulation', 'симуля', 'spice', 'oscilloscope', 'осцилл'),
        'reliability': ('reliability', 'надеж', 'надёж', 'mtbf', 'failure'),
        'manufacturing': ('assembly', 'сборк', 'монтаж', 'пайк', 'технолог'),
        'economics': ('cost', 'цена', 'стоим', 'себестоим', 'bom'),
        'safety': ('safety', 'безопас', 'заземл', 'ground'),
        'fault': ('fault', 'ошиб', 'дефект', 'отказ', 'short'),
    }
    return [key for key, needles in candidates.items() if any(needle in lowered for needle in needles)]


def _process_steps(text: str) -> list[dict[str, str]]:
    steps = []
    for key, label in [
        ('requirements', 'requirements'),
        ('schematic', 'schematic design'),
        ('simulation', 'simulation or measurement'),
        ('reliability', 'reliability calculation'),
        ('manufacturing', 'manufacturing readiness'),
        ('economics', 'economic assessment'),
        ('safety', 'safety check'),
        ('fault', 'fault diagnostics'),
    ]:
        if key in _keywords(text):
            steps.append({'kind': key, 'label': label})
    return steps


def _first_int(pattern: str, text: str) -> int | None:
    match = re.search(pattern, text, flags=re.I)
    if not match:
        return None
    try:
        return int(match.group(1))
    except TypeError, ValueError:
        return None
