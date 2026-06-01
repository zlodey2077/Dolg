from __future__ import annotations

import shutil
from pathlib import Path

from docx import Document
from docx.shared import Pt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
SCREENSHOTS = DOCS / "diploma_assets" / "screenshots"
P6 = SCREENSHOTS / "presentation_v6"

TITLE = (
    "Разработка веб-приложения для продажи радио- и электронных компонентов "
    "со встроенными инструментами проектирования и симуляции схем"
)
AUTHOR = "Буряко Дмитрий Сергеевич"
SUPERVISOR = "Буланов Сергей Георгиевич"

OUT_PPTX = DOCS / "Презентация_DOLG_основная_защита_20_слайдов_20260524.pptx"
LEGACY_CURRENT_PPTX = DOCS / "Презентация_DOLG_актуальная_20260519.pptx"
OUT_MD = DOCS / "Речь_и_вопросы_к_защите_DOLG_20260513_v5.md"
OUT_DOCX = DOCS / "Речь_и_вопросы_к_защите_DOLG_20260524_20_слайдов.docx"
LEGACY_SPEECH_DOCX = DOCS / "Речь_и_вопросы_к_защите_DOLG_20260513_v5_с_вопросами_по_схеме_20260519.docx"


class C:
    bg = RGBColor(220, 228, 237)
    bg2 = RGBColor(207, 218, 230)
    panel = RGBColor(238, 243, 248)
    panel2 = RGBColor(226, 234, 242)
    navy = RGBColor(15, 23, 42)
    text = RGBColor(30, 41, 59)
    muted = RGBColor(71, 85, 105)
    line = RGBColor(148, 163, 184)
    cyan = RGBColor(8, 145, 178)
    blue = RGBColor(37, 99, 235)
    green = RGBColor(22, 163, 74)
    orange = RGBColor(234, 88, 12)
    purple = RGBColor(124, 58, 237)
    red = RGBColor(220, 38, 38)
    code_bg = RGBColor(232, 238, 245)
    header = RGBColor(214, 224, 236)
    rail = RGBColor(194, 207, 222)


def e(value: float):
    return Inches(value)


def set_run(run, size: float, color: RGBColor, *, bold: bool = False, font: str = "Arial") -> None:
    run.font.name = font
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = color


class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = e(16)
        self.prs.slide_height = e(9)
        self.total = 20

    def slide(self, number: int, title: str = ""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = C.bg
        rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, e(0.18), e(0.34), e(8.82))
        rail.fill.solid()
        rail.fill.fore_color.rgb = C.rail
        rail.line.fill.background()
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, e(16), e(0.18))
        band.fill.solid()
        band.fill.fore_color.rgb = C.cyan
        band.line.fill.background()
        self.text(slide, f"{number}/{self.total}", 15.15, 8.45, 0.55, 0.25, 13, C.muted, align="right")
        if title:
            self.block(slide, 0.58, 0.34, 13.3, 1.0, fill=C.header)
            self.text(slide, title, 0.8, 0.47, 12.85, 0.72, 38, C.navy, bold=True)
        return slide

    def text(
        self,
        slide,
        value: str,
        x: float,
        y: float,
        w: float,
        h: float,
        size: float,
        color: RGBColor = C.text,
        *,
        bold: bool = False,
        align: str = "left",
        font: str = "Arial",
        line_spacing: float | None = None,
    ):
        shape = slide.shapes.add_textbox(e(x), e(y), e(w), e(h))
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        frame.vertical_anchor = MSO_ANCHOR.TOP
        frame.margin_left = e(0.01)
        frame.margin_right = e(0.01)
        frame.margin_top = e(0.01)
        frame.margin_bottom = e(0.01)
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        for i, line in enumerate(value.splitlines() or [""]):
            p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
            p.alignment = align_map[align]
            if line_spacing is not None:
                p.line_spacing = line_spacing
            run = p.add_run()
            run.text = line
            set_run(run, size, color, bold=bold, font=font)
        return shape

    def block(self, slide, x: float, y: float, w: float, h: float, *, fill: RGBColor = C.panel, accent: RGBColor | None = None):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, e(x), e(y), e(w), e(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = C.line
        shape.line.width = PptPt(1.25)
        if accent:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, e(x), e(y), e(w), e(0.12))
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent
            bar.line.fill.background()
        return shape

    def tile(self, slide, x: float, y: float, w: float, h: float, title: str, body: str, accent: RGBColor, *, body_size: float = 24) -> None:
        self.block(slide, x, y, w, h, accent=accent)
        title_size = 24 if len(title) > 16 else 26
        self.text(slide, title, x + 0.18, y + 0.18, w - 0.36, 0.76, title_size, C.navy, bold=True)
        self.text(slide, body, x + 0.18, y + 0.86, w - 0.36, h - 0.92, body_size, C.text, line_spacing=0.92)

    def bullet_tile(self, slide, x: float, y: float, w: float, h: float, title: str, items: list[str], accent: RGBColor, *, size: float = 26) -> None:
        if len(items) >= 6:
            size = min(size, 21)
        elif len(items) >= 5:
            size = min(size, 22)
        if h < 3.0:
            size = min(size, 22)
        if h < 2.4:
            size = min(size, 20)
        body = "\n".join(f"• {item}" for item in items)
        self.tile(slide, x, y, w, h, title, body, accent, body_size=size)

    def mini_tile(self, slide, x: float, y: float, w: float, h: float, title: str, body: str, accent: RGBColor) -> None:
        self.block(slide, x, y, w, h, accent=accent)
        lines = max(1, len(body.splitlines()))
        body_size = 16.5 if len(body) > 7 else (17.5 if lines >= 3 else 19)
        self.text(slide, title, x + 0.12, y + 0.2, w - 0.24, 0.4, 22, C.navy, bold=True, align="center")
        self.text(slide, body, x + 0.12, y + 0.6, w - 0.24, h - 0.66, body_size, C.text, align="center")

    def metric(self, slide, x: float, y: float, w: float, h: float, value: str, label: str, accent: RGBColor) -> None:
        self.block(slide, x, y, w, h, accent=accent, fill=C.panel)
        self.text(slide, value, x + 0.1, y + 0.21, w - 0.2, 0.44, 30, accent, bold=True, align="center")
        self.text(slide, label, x + 0.08, y + 0.78, w - 0.16, 0.45, 18.5, C.text, bold=True, align="center")

    def image_full(self, slide, path: Path, x: float, y: float, w: float, h: float) -> None:
        if not path.exists():
            self.text(slide, f"Нет файла: {path.name}", x, y + h / 2, w, 0.5, 28, C.red, bold=True, align="center")
            return
        with Image.open(path) as im:
            iw, ih = im.size
        scale = min(w / iw, h / ih)
        pw = iw * scale
        ph = ih * scale
        px = x + (w - pw) / 2
        py = y + (h - ph) / 2
        slide.shapes.add_picture(str(path), e(px), e(py), width=e(pw), height=e(ph))

    def table(self, slide, x: float, y: float, w: float, h: float, headers: list[str], rows: list[tuple[str, ...]], widths: list[float]) -> None:
        shape = slide.shapes.add_table(len(rows) + 1, len(headers), e(x), e(y), e(w), e(h))
        table = shape.table
        for i, width in enumerate(widths):
            table.columns[i].width = e(width)
        for c, header in enumerate(headers):
            cell = table.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C.navy
            cell.text = header
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    set_run(r, 22, RGBColor(255, 255, 255), bold=True)
        for r_i, row in enumerate(rows, start=1):
            for c_i, value in enumerate(row):
                cell = table.cell(r_i, c_i)
                cell.fill.solid()
                cell.fill.fore_color.rgb = C.panel if r_i % 2 else C.panel2
                cell.text = value
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.LEFT
                    for r in p.runs:
                        set_run(r, 20, C.text)

    def code_plain(self, slide, value: str, x: float, y: float, w: float, h: float, *, size: float = 20) -> None:
        self.block(slide, x, y, w, h, fill=C.code_bg)
        self.text(slide, value, x + 0.18, y + 0.16, w - 0.36, h - 0.26, size, C.navy, font="Consolas", line_spacing=0.88)

    def arrow_flow(self, slide, items: list[str], x: float, y: float, w: float, h: float, accent: RGBColor = C.cyan, *, size: float = 21) -> None:
        cell_w = w / len(items)
        for i, item in enumerate(items):
            sx = x + i * cell_w
            self.block(slide, sx, y, cell_w - 0.13, h, fill=C.panel, accent=accent if i == 0 else None)
            self.text(slide, item, sx + 0.08, y + 0.38, cell_w - 0.28, h - 0.45, size, C.navy, bold=True, align="center")
            if i < len(items) - 1:
                self.text(slide, "→", sx + cell_w - 0.16, y + 0.43, 0.16, 0.25, 22, accent, bold=True, align="center")


def build_deck() -> None:
    d = Deck()

    s = d.slide(1)
    d.block(s, 0.72, 0.65, 6.95, 4.95, fill=C.header, accent=C.cyan)
    d.block(s, 0.72, 5.95, 6.95, 1.75, fill=C.panel, accent=C.green)
    d.text(s, "Выпускная квалификационная работа", 0.95, 0.88, 6.45, 0.38, 23, C.cyan, bold=True)
    d.text(s, TITLE, 0.95, 1.52, 6.22, 3.35, 27, C.navy, bold=True, line_spacing=0.86)
    d.text(s, f"Автор: {AUTHOR}\nРуководитель: {SUPERVISOR}\nГод защиты: 2026", 0.95, 6.22, 6.25, 1.05, 18.5, C.text, bold=True, line_spacing=0.9)
    d.image_full(s, SCREENSHOTS / "01_home_catalog.png", 8.2, 1.25, 7.1, 5.55)

    s = d.slide(2, "Проблемы, цели и задачи")
    d.tile(s, 0.75, 1.65, 4.45, 5.35, "Проблемы", "Каталог, datasheet, CAD, SPICE, BOM и заказ часто существуют отдельно.\n\nНоминалы, модели и результаты приходится переносить вручную.", C.blue, body_size=25)
    d.tile(s, 5.78, 1.65, 4.45, 5.35, "Цель", "Разработать веб-приложение для продажи радио- и электронных компонентов со встроенными инструментами проектирования и симуляции схем.", C.orange, body_size=25)
    d.tile(s, 10.82, 1.65, 4.45, 5.35, "Задачи", "Каталог и заказы.\nCAD/SIM.\nЛаборатория.\nОбучение.\nReview.\nКачество данных.", C.green, body_size=25)

    s = d.slide(3, "Анализ решений")
    d.table(
        s,
        0.7,
        1.55,
        14.6,
        5.75,
        ["Класс", "Что хорошо", "Чего не хватает"],
        [
            ("Маркетплейсы", "поиск, цена, заказ", "нет проверки схемы"),
            ("Онлайн-симуляторы", "быстрый эксперимент", "слабая связь с товарами"),
            ("KiCad / Altium", "EDA workflow, DRC/ERC", "высокий порог входа"),
            ("AI-CAD", "подсказки по проекту", "нужны проверяемые факты"),
        ],
        [3.4, 5.4, 5.8],
    )

    s = d.slide(4, "Функции системы")
    functions = [
        ("Каталог", "карточки\nфильтры\nпараметры", C.blue),
        ("Knowledge", "статьи\nматериалы\nобучение", C.purple),
        ("Лаборатория", "расчеты\nоценки\nметрики", C.green),
        ("CAD", "схема\nGND\nэкспорт", C.cyan),
        ("SIM", "DC / AC / TRAN\nFFT / Bode\nMonte Carlo", C.orange),
        ("Review", "DRC/ERC\nBOM risk\nfindings", C.red),
        ("Import", "LTspice\nSPICE\nKiCad subset", C.blue),
        ("Заказ", "BOM\nXLSX\nкорзина", C.green),
        ("AI", "правила\nфакты\nподсказки", C.purple),
    ]
    for idx, (title, body, color) in enumerate(functions):
        x = 0.7 + (idx % 3) * 5.05
        y = 1.45 + (idx // 3) * 2.15
        d.mini_tile(s, x, y, 4.45, 1.65, title, body, color)

    s = d.slide(5, "Архитектура: приложения и данные")
    d.bullet_tile(s, 0.75, 1.45, 4.65, 5.8, "Django apps", ["shop", "accounts", "orders", "knowledge", "Dolg_APP"], C.blue, size=30)
    d.bullet_tile(s, 5.72, 1.45, 4.65, 5.8, "Данные", ["товары", "схемы", "измерения", "уроки", "review snapshots", "заказы"], C.green, size=28)
    d.bullet_tile(s, 10.7, 1.45, 4.55, 5.8, "Контуры", ["каталог", "проектирование", "симуляция", "обучение", "коммерция"], C.orange, size=29)

    s = d.slide(6, "Архитектура: service-layer и библиотеки")
    d.arrow_flow(s, ["UI", "Services", "Libraries", "Data", "Reports"], 0.8, 1.55, 14.4, 1.05, size=25)
    d.bullet_tile(s, 0.8, 3.05, 4.3, 3.4, "Scientific", ["NumPy", "SciPy", "Matplotlib", "Pandas"], C.blue, size=26)
    d.bullet_tile(s, 5.85, 3.05, 4.3, 3.4, "Graph / Formula", ["NetworkX", "SymPy", "Schemdraw", "Pint"], C.green, size=26)
    d.bullet_tile(s, 10.9, 3.05, 4.3, 3.4, "Expert", ["jsonschema", "rule-engine", "Lark", "Z3", "scikit-fuzzy"], C.orange, size=25)

    s = d.slide(7, "Каталог")
    d.image_full(s, P6 / "01_catalog_cards_current.png", 0.55, 1.28, 9.25, 6.55)
    d.bullet_tile(s, 10.12, 1.35, 5.15, 5.95, "Функции", ["карточки товаров", "активные chips", "unit-aware фильтры", "datasheet / SPICE / CAD", "media quality 89/89"], C.cyan, size=26)

    s = d.slide(8, "Knowledge и обучение")
    d.image_full(s, SCREENSHOTS / "08_knowledge_article_materials.png", 0.55, 1.35, 8.4, 5.95)
    d.mini_tile(s, 9.35, 1.55, 1.75, 1.15, "21", "статья", C.blue)
    d.mini_tile(s, 11.35, 1.55, 1.75, 1.15, "50", "материалы", C.cyan)
    d.mini_tile(s, 13.35, 1.55, 1.75, 1.15, "29", "задания", C.green)
    d.bullet_tile(s, 9.35, 3.25, 5.8, 3.1, "Типы заданий", ["числовой ответ", "сборка схемы", "измерение симуляции"], C.orange, size=28)

    s = d.slide(9, "Инженерная лаборатория")
    d.image_full(s, P6 / "02_engineering_lab_results.png", 0.55, 1.15, 9.55, 6.75)
    d.bullet_tile(s, 10.45, 1.35, 4.8, 3.35, "Расчеты", ["NE555", "стабилизатор", "транзисторный ключ", "тепловой запас", "RC-антидребезг"], C.blue, size=22)
    d.bullet_tile(s, 10.45, 5.05, 4.8, 2.35, "Оценки", ["норма", "риск", "перегрев", "нужен запас"], C.orange, size=23)

    s = d.slide(10, "CAD-редактор")
    d.image_full(s, P6 / "03_cad_editor_current.png", 0.55, 1.15, 11.1, 6.8)
    d.bullet_tile(s, 12.0, 1.35, 3.2, 5.75, "Режим", ["GND", "источник", "components", "wires", "labels", "export"], C.cyan, size=26)

    s = d.slide(11, "Симуляция и измерения")
    d.image_full(s, P6 / "04_simulation_current.png", 0.55, 1.15, 9.5, 6.7)
    d.bullet_tile(s, 10.4, 1.25, 4.8, 2.55, "Анализы", ["DC / OP", "AC", "TRAN"], C.blue, size=24)
    d.bullet_tile(s, 10.4, 4.15, 4.8, 3.35, "Метрики", ["напряжение узла", "ток ветви", "RMS / частота / duty", "мощность / температура", "expected vs measured"], C.green, size=21)

    s = d.slide(12, "Pro-аналитика")
    d.image_full(s, SCREENSHOTS / "02_simulation_ac_graph.png", 0.55, 1.15, 9.6, 6.7)
    d.bullet_tile(s, 10.55, 1.45, 4.65, 2.85, "Возможности", ["FFT", "Bode plot", "Monte Carlo", "parameter sweep"], C.orange, size=23)
    d.bullet_tile(s, 10.55, 4.75, 4.65, 2.55, "Стек", ["NumPy / SciPy", "Matplotlib / Pandas", "scientific stack OK"], C.green, size=21)

    s = d.slide(13, "Экспертные уровни контроля")
    d.arrow_flow(s, ["Правила", "Граф", "Формулы", "Расчет", "AI-подсказка", "Человек"], 0.65, 1.55, 14.7, 1.2, C.purple, size=23)
    d.tile(s, 1.05, 3.35, 4.2, 2.9, "Проверяемые факты", "GND, источник, номиналы, связи, BOM, измерения.", C.blue, body_size=23)
    d.tile(s, 5.9, 3.35, 4.2, 2.9, "Автоматический вывод", "findings, риски, рекомендации, подбор вариантов.", C.orange, body_size=23)
    d.tile(s, 10.75, 3.35, 4.2, 2.9, "Финальный контроль", "Последнее решение всегда остается за человеком.", C.green, body_size=23)

    s = d.slide(14, "CAD-импорт и базы")
    d.arrow_flow(s, ["LTspice", "SPICE", "KiCad subset", "scheme_data", "review"], 0.8, 1.55, 14.4, 1.1, C.cyan, size=23)
    bases = [
        ("Компоненты", "параметры\nналичие\nкарточки", C.blue),
        ("SPICE / CAD", "модели\nкорпуса\nnetlist", C.green),
        ("Правила", "DRC/ERC\nrisk\nfindings", C.orange),
        ("Знания", "статьи\nматериалы\nформулы", C.purple),
        ("Задания", "math\ncircuit\nmeasure", C.red),
    ]
    for idx, (title, body, color) in enumerate(bases):
        d.mini_tile(s, 0.8 + idx * 2.92, 3.55, 2.45, 2.6, title, body, color)

    s = d.slide(15, "BOM и заказ")
    d.image_full(s, SCREENSHOTS / "07_product_related_cards.png", 0.65, 1.4, 7.15, 3.45)
    d.bullet_tile(s, 8.15, 1.35, 3.25, 4.8, "BOM", ["позиции схемы", "подбор товара", "XLSX"], C.green, size=27)
    d.bullet_tile(s, 11.85, 1.35, 3.25, 4.8, "Заказ", ["корзина", "остатки", "checkout"], C.orange, size=27)
    d.metric(s, 0.95, 5.55, 3.0, 1.35, "3/3", "browser smoke", C.blue)
    d.metric(s, 4.35, 5.55, 3.0, 1.35, "OK", "order flow", C.green)

    s = d.slide(16, "Развитие проекта")
    d.bullet_tile(s, 0.8, 1.55, 6.8, 5.65, "Ближайшее", ["комментарии пользователей", "measurement core", "official image/API sources", "datasheet extraction"], C.blue, size=29)
    d.bullet_tile(s, 8.4, 1.55, 6.8, 5.65, "Дальнейшее", ["OR-Tools для BOM", "RDFLib для онтологии", "neural deep hints", "расширение CAD-import"], C.purple, size=29)

    s = d.slide(17, "Код: сервисы")
    d.code_plain(s, "def build_design_review(project):\n    scheme = project.scheme_data\n    validation = validate_scheme_data(scheme)\n    metrics = build_connectivity_metrics(scheme)\n    return compose_review(validation, metrics)", 0.75, 1.35, 6.95, 2.55, size=17)
    d.code_plain(s, "def grade_task(task, payload):\n    if task.task_type == 'math_numeric':\n        answer = payload['answer']\n        return grade_math_task(task, answer)\n    scheme = payload['scheme_data']\n    return grade_circuit_task(task, scheme)", 8.15, 1.35, 7.05, 2.85, size=16)
    d.code_plain(s, "parsed = parse_engineering_quantity('10 кОм', expected_unit='ohm')\n# 10000.0 ohm", 3.0, 4.55, 10.0, 1.3, size=20)
    d.block(s, 1.05, 6.3, 13.7, 0.98, fill=C.header, accent=C.green)
    d.text(s, "Код вынесен в сервисы: один слой работает в лаборатории, обучении и review.", 1.35, 6.55, 13.1, 0.38, 26, C.navy, bold=True, align="center")

    s = d.slide(18, "Код: scheme_data и finding")
    d.code_plain(s, "scheme_data = {\n  'components': [\n    {'id':'V1', 'type':'source', 'value':'9 V'},\n    {'id':'R1', 'type':'resistor', 'value':'20 kOhm'},\n    {'id':'R2', 'type':'resistor', 'value':'10 kOhm'}\n  ],\n  'labels': {'output':'vout'}\n}", 0.75, 1.3, 7.0, 4.7, size=15.5)
    d.code_plain(s, "finding = {\n  'rule_id': 'circuit.ground.required',\n  'severity': 'error',\n  'evidence': ['GND node not found'],\n  'recommendation': 'Добавьте опорную землю'\n}", 8.25, 1.3, 7.0, 4.7, size=16.5)
    d.block(s, 1.0, 6.34, 14.0, 0.98, fill=C.header, accent=C.orange)
    d.text(s, "Схема, ошибка и рекомендация представлены как данные.\nИх можно проверять, сохранять и объяснять.", 1.3, 6.5, 13.4, 0.62, 22, C.navy, bold=True, align="center", line_spacing=0.9)

    s = d.slide(19, "Итоги")
    d.metric(s, 0.8, 1.45, 2.3, 1.45, "89", "товаров", C.blue)
    d.metric(s, 3.45, 1.45, 2.3, 1.45, "43", "РЭБ", C.cyan)
    d.metric(s, 6.1, 1.45, 2.3, 1.45, "21", "статья", C.green)
    d.metric(s, 8.75, 1.45, 2.3, 1.45, "29", "заданий", C.orange)
    d.metric(s, 11.4, 1.45, 3.2, 1.45, "OK", "demo-ready", C.purple)
    d.bullet_tile(s, 1.25, 3.55, 6.55, 3.55, "Создано", ["каталог + заказы", "CAD/SIM", "лаборатория", "обучение", "expert review"], C.blue, size=24)
    d.bullet_tile(s, 8.45, 3.55, 6.25, 3.55, "Проверено", ["manage.py check", "check_demo_ready", "check_data_integrity", "media quality 89/89"], C.green, size=24)

    s = d.slide(20, "Спасибо за внимание")
    d.block(s, 0.72, 1.32, 9.0, 2.2, fill=C.header, accent=C.cyan)
    d.block(s, 0.72, 4.05, 9.45, 2.85, fill=C.panel, accent=C.green)
    d.text(s, "DOLG", 0.9, 1.6, 5.2, 0.9, 58, C.cyan, bold=True)
    d.text(s, "каталог • CAD/SIM • лаборатория\nобучение • review", 0.95, 2.55, 8.4, 0.9, 25, C.navy, bold=True, line_spacing=0.9)
    d.text(s, f"{AUTHOR}\nНаучный руководитель: {SUPERVISOR}\nEmail: buryako@internet.com\nТелефон: +7 (903) 439-44-87", 0.95, 4.36, 9.05, 2.1, 23.5, C.text, bold=True)
    d.image_full(s, SCREENSHOTS / "01_home_catalog.png", 10.2, 3.85, 4.85, 2.95)

    d.prs.save(OUT_PPTX)
    shutil.copy2(OUT_PPTX, LEGACY_CURRENT_PPTX)


SPEECH_SECTIONS = [
    ("Слайд 1. Титул", "Добрый день. Тема моей выпускной квалификационной работы - «Разработка веб-приложения для продажи радио- и электронных компонентов со встроенными инструментами проектирования и симуляции схем». Сегодня я покажу проект как единую систему: каталог, проектирование, симуляция, обучение, экспертная проверка и заказ."),
    ("Слайд 2. Проблемы, цели и задачи", "Основная проблема в том, что инженерный процесс часто разорван между каталогом, datasheet, CAD, SPICE и BOM. Из-за ручного переноса данных легко ошибиться в номинале, единицах измерения или модели. Цель работы - собрать эти действия в одном веб-приложении. Задачи закрывают каталог, учетные записи, корзину и заказы, CAD/SIM, лабораторию, обучение и контроль качества данных."),
    ("Слайд 3. Анализ решений", "Существующие решения закрывают отдельные части задачи. Маркетплейсы хорошо продают и ищут товар, но не анализируют схему. Онлайн-симуляторы удобны для экспериментов, но редко связаны с реальными карточками компонентов. KiCad и Altium сильны как промышленные CAD-системы, но сложны для учебного маршрута. AI-CAD полезен, но инженерный вывод должен оставаться объяснимым."),
    ("Слайд 4. Функции системы", "На этом слайде показаны основные функции DOLG. Каталог и фильтры отвечают за подбор компонентов. Knowledge и обучение дают теорию и практические задания. Лаборатория выполняет расчеты, CAD и SIM позволяют собрать и проверить схему. Review и AI-помощник объясняют ошибки, а BOM и заказ связывают инженерную часть с покупкой."),
    ("Слайд 5. Архитектура: приложения и данные", "Серверная часть построена на Django-приложениях. Shop отвечает за каталог, accounts - за пользователей, orders - за заказы, knowledge - за статьи и обучение, Dolg_APP - за CAD/SIM и инженерные сервисы. Данные включают товары, схемы, измерения, уроки, отчеты review и заказы."),
    ("Слайд 6. Архитектура: service-layer и библиотеки", "Ключевая идея архитектуры - service-layer. UI только показывает результат, а расчеты, проверки, импорт, обучение и AI-помощник используют общие сервисы. Scientific stack дает FFT, Bode и Monte Carlo. Graph/formula stack анализирует схему и формулы. Expert stack отвечает за правила, единицы измерения, импорт и подбор вариантов."),
    ("Слайд 7. Каталог", "Каталог показывает карточки товаров, активные chips, инженерные параметры и flags datasheet, SPICE и CAD. Поиск поддерживает единицы измерения и диапазоны. Качество изображений контролируется отдельно: активные карточки не должны получать случайные нерелевантные фотографии."),
    ("Слайд 8. Knowledge и обучение", "Knowledge объединяет статьи, материалы и практикум. Обучение не сводится к тестам с выбором ответа: есть числовые задачи, задачи на сборку схемы и задания на измерение результата симуляции. Это позволяет учиться через реальные инженерные сценарии."),
    ("Слайд 9. Инженерная лаборатория", "Лаборатория содержит прикладные расчеты: NE555, стабилизатор, транзисторный ключ, тепловой запас и RC-антидребезг. Важный момент - результат сопровождается инженерной оценкой. Пользователь видит не только число, но и статус: норма, риск, перегрев или нужен запас."),
    ("Слайд 10. CAD-редактор", "CAD-редактор нужен не только для рисования. Схема содержит GND, источник, компоненты, провода и labels. Затем эта схема становится входом для DRC/ERC, netlist, BOM, симуляции и Engineering Review."),
    ("Слайд 11. Симуляция и измерения", "Симулятор поддерживает DC, AC и TRAN-анализы. Результат можно оценивать по метрикам: напряжение узла, ток ветви, RMS, частота, duty cycle, мощность и температура. Отдельный сценарий - expected vs measured, когда расчетное значение сравнивается с измеренным."),
    ("Слайд 12. Pro-аналитика", "Pro-аналитика расширяет лабораторию и симулятор. NumPy и SciPy используются для численных расчетов, Matplotlib - для графиков, Pandas - для агрегированной статистики. Это позволяет строить FFT, Bode plot, Monte Carlo tolerance и sweep по параметрам."),
    ("Слайд 13. Экспертные уровни контроля", "Expert-first подход означает, что сначала работают проверяемые правила, граф схемы, формулы и единицы измерения. Затем подключаются расчет, оптимизация и AI-подсказка. Но последний уровень контроля всегда остается за человеком: система помогает увидеть риск, а инженер принимает финальное решение."),
    ("Слайд 14. CAD-импорт и базы", "Импорт нужен для того, чтобы простую схему из LTspice, SPICE или KiCad subset привести к внутреннему scheme_data и прогнать через review. Проект также опирается на базы компонентов, моделей, правил, знаний и учебных задач."),
    ("Слайд 15. BOM и заказ", "Коммерческий контур включает BOM, подбор товара из каталога, XLSX-экспорт, корзину, учет остатков и оформление заказа. Так результат проектирования превращается в набор компонентов для покупки."),
    ("Слайд 16. Развитие проекта", "Ближайшее развитие - комментарии пользователей, measurement core, официальные источники изображений и API, а также расширение извлечения данных из datasheet. Дальнейшее развитие - OR-Tools для оптимизации BOM, RDFLib для онтологии компонентов, neural deep hints и расширенный CAD-import."),
    ("Слайд 17. Код: сервисы", "Кодовые фрагменты показывают, что расчеты и проверки вынесены в сервисы. Review строится через валидацию схемы и связность, grader проверяет учебные задания, unit-service приводит номиналы к единым единицам. На защите этот слайд нужен как короткое доказательство реализации."),
    ("Слайд 18. Код: scheme_data и finding", "Схема и результат проверки представлены как данные. Scheme_data хранит компоненты и связи, а finding хранит правило, серьезность, доказательства и рекомендацию. Поэтому отчет можно сохранять, объяснять и использовать в обучении."),
    ("Слайд 19. Итоги", "В результате создана платформа, где каталог, CAD/SIM, лаборатория, обучение, review и заказ работают в одном контуре. Проект имеет проверяемые данные и проходит основные проверки: manage.py check, demo-ready, data-integrity и media quality."),
    ("Слайд 20. Спасибо за внимание", "Спасибо за внимание. Готов ответить на вопросы по архитектуре, демонстрационной схеме, экспертным правилам, симуляции и дальнейшему развитию проекта."),
]

SCHEME_QUESTIONS = [
    ("Почему для демонстрации выбран делитель напряжения?", "Это простая схема с источником, GND, двумя номиналами, выходным узлом и измеряемым результатом. На ней видно весь маршрут DOLG."),
    ("Какая формула используется?", "Vout = Vin * R2 / (R1 + R2), где R1 подключен к источнику, R2 к земле, а Vout находится между ними."),
    ("Как получить около 3 В из 9 В?", "Нужно отношение R2 / (R1 + R2) около 1/3. Например, R1 = 20 кОм и R2 = 10 кОм."),
    ("Почему GND критичен?", "Без GND нет опорного потенциала. Симулятор и review не могут корректно трактовать напряжения узлов."),
    ("Что проверяет Engineering Review?", "Наличие GND и источника, связность графа, номиналы, output node, BOM-связи и совпадение expected с measured."),
    ("Что означает expected vs measured?", "Expected - расчетное значение по формуле, measured - результат симуляции или измерения."),
    ("Какие ошибки удобно показать?", "Нет GND, перепутанные резисторы, неверный узел измерения, слишком малые сопротивления и перегрузка по мощности."),
    ("Как нагрузка влияет на делитель?", "Нагрузка подключается параллельно нижнему плечу и меняет выходное напряжение."),
    ("Почему важна мощность резисторов?", "Правильное напряжение не гарантирует безопасный режим. Нужно проверять запас по мощности."),
    ("Как схема становится заданием?", "Пользователь считает Vout, собирает схему, запускает DC-анализ и отправляет измерение в grader."),
    ("Что меняется для RC-фильтра?", "Появляется частотная область: расчет fc, Bode plot и точка около -3 дБ."),
    ("Что можно спросить по LED-ветви?", "Как выбрать резистор, ток светодиода, падение напряжения и запас по мощности."),
    ("Это калькулятор или инженерная проверка?", "Это инженерная проверка, потому что формула связана со схемой, графом, симуляцией, BOM и finding."),
    ("Почему финальный контроль за человеком?", "Автоматические правила находят риски и дают рекомендации, но инженер учитывает контекст, допущения и реальные условия эксплуатации."),
]

GENERAL_QA = [
    ("Почему expert-first, а не сразу нейронная сеть?", "Инженерный вывод должен быть объяснимым. Нейронная сеть может появиться позже как подсказка, но не как финальный verdict."),
    ("Чем DOLG отличается от KiCad или Altium?", "DOLG не заменяет промышленный CAD. Он связывает компонент, расчет, схему, симуляцию, обучение и заказ."),
    ("Зачем нужен Pint?", "Он приводит пользовательские значения вроде 10k, 10 кОм и 100 нФ к единому виду."),
    ("Зачем нужен Z3?", "Он помогает подбирать допустимые номиналы при ограничениях."),
    ("Что дает Datasheet Intelligence?", "Это база для извлечения pinout, package, ratings и thermal data из datasheet."),
    ("Почему важна media-policy?", "Карточки не должны получать случайные нерелевантные изображения. Качество медиа проверяется автоматически."),
    ("Какие проверки подтверждают готовность?", "manage.py check, check_demo_ready, check_data_integrity, targeted tests и browser smoke."),
    ("Что развивать дальше?", "Комментарии, measurement core, official sources, OR-Tools, RDFLib и neural deep hints после expert baseline."),
]


def build_speech_text() -> str:
    lines = [
        "# Речь к основной защите DOLG",
        "",
        "Речь синхронизирована с новой крупной презентацией на 20 слайдов. На слайдах оставлена только основная информация, подробности раскрываются устно.",
        "",
        "## Речь по слайдам",
        "",
    ]
    for title, body in SPEECH_SECTIONS:
        lines.extend([f"### {title}", body, ""])
    lines.extend([
        "## Вопросы по разбираемой схеме",
        "",
        "Основная демонстрационная схема - делитель напряжения 9 В -> около 3 В с узлом Vout, GND и двумя резисторами.",
        "",
    ])
    for i, (question, answer) in enumerate(SCHEME_QUESTIONS, start=1):
        lines.extend([f"{i}. **{question}**", f"   {answer}", ""])
    lines.extend(["## Возможные вопросы и ответы", ""])
    for i, (question, answer) in enumerate(GENERAL_QA, start=1):
        lines.extend([f"{i}. **{question}**", f"   {answer}", ""])
    return "\n".join(lines).rstrip() + "\n"


def build_speech_docx(markdown_text: str) -> None:
    doc = Document()
    doc.styles["Normal"].font.name = "Times New Roman"
    doc.styles["Normal"].font.size = Pt(12)
    for line in markdown_text.splitlines():
        if line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.strip():
            p = doc.add_paragraph()
            run = p.add_run(line.replace("**", ""))
            run.font.name = "Times New Roman"
            run.font.size = Pt(12)
        else:
            doc.add_paragraph()
    doc.save(OUT_DOCX)
    shutil.copy2(OUT_DOCX, LEGACY_SPEECH_DOCX)


def main() -> None:
    build_deck()
    speech = build_speech_text()
    OUT_MD.write_text(speech, encoding="utf-8")
    build_speech_docx(speech)
    print(f"PPTX: {OUT_PPTX}")
    print(f"PPTX copy: {LEGACY_CURRENT_PPTX}")
    print(f"Speech MD: {OUT_MD}")
    print(f"Speech DOCX: {OUT_DOCX}")
    print(f"Speech DOCX copy: {LEGACY_SPEECH_DOCX}")


if __name__ == "__main__":
    main()
