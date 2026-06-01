from __future__ import annotations

import shutil
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt as DocPt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
DOWNLOADS = Path.home() / "Downloads"

TITLE = "Разработка веб-приложения для продажи радио- и электронных компонентов со встроенными инструментами проектирования и симуляции схем"
AUTHOR = "Буряко Дмитрий Сергеевич"
SUPERVISOR = "Буланов Сергей Георгиевич"


def first_file(pattern: str) -> Path:
    files = [path for path in DOCS.glob(pattern) if not path.name.startswith("~$")]
    if not files:
        raise FileNotFoundError(pattern)
    return files[0]


def backup_file(path: Path, backup_dir: Path) -> None:
    if path.exists():
        target = backup_dir / path.name
        if not target.exists():
            shutil.copy2(path, target)


def unlocked_variant(path: Path) -> Path:
    return path.with_name(f"{path.stem}_актуальная_20260519{path.suffix}")


def save_docx_safe(doc: Document, path: Path) -> Path:
    try:
        doc.save(str(path))
        return path
    except PermissionError:
        fallback = unlocked_variant(path)
        if fallback == path:
            fallback = path.with_name(f"{path.stem}_copy{path.suffix}")
        doc.save(str(fallback))
        return fallback


def save_pptx_safe(prs: Presentation, path: Path) -> Path:
    try:
        prs.save(str(path))
        return path
    except PermissionError:
        fallback = unlocked_variant(path)
        if fallback == path:
            fallback = path.with_name(f"{path.stem}_copy{path.suffix}")
        prs.save(str(fallback))
        return fallback


def copy_file_safe(source: Path, target: Path) -> Path:
    try:
        shutil.copy2(source, target)
        return target
    except PermissionError:
        fallback = unlocked_variant(target)
        if fallback == target:
            fallback = target.with_name(f"{target.stem}_copy{target.suffix}")
        shutil.copy2(source, fallback)
        return fallback


def add_run(paragraph, text: str, *, size: int = 12, bold: bool = False) -> None:
    run = paragraph.add_run(text)
    run.font.name = "Times New Roman"
    run.font.size = DocPt(size)
    run.font.bold = bold


def replace_paragraph(paragraph, text: str) -> None:
    paragraph.clear()
    add_run(paragraph, text)


def update_diploma() -> None:
    diploma_path = first_file("Диплом*.docx")
    doc = Document(str(diploma_path))

    replacements = {
        "Выпускная квалификационная работа достигла поставленной цели:": (
            "Выпускная квалификационная работа достигла поставленной цели: разработана демонстрационная веб-платформа DOLG, "
            "объединяющая каталог радио- и электронных компонентов, покупательский контур, редактор схем, SPICE-симуляцию, "
            "инженерную лабораторию, практикум обучения и Engineering Review проекта. Реализация выполнена на Django с "
            "использованием ORM, шаблонов, встроенной аутентификации, локальных статических библиотек и переиспользуемого "
            "service-layer для расчетов, проверок, импорта и отчетов."
        ),
        "Практическая значимость проекта состоит": (
            "Практическая значимость проекта состоит в объединении действий, которые обычно выполняются в разрозненных системах: "
            "подбор компонента, построение схемы, расчет, симуляция, измерение, генерация BOM, экспертная проверка и подготовка заказа. "
            "В актуальной версии проекта зафиксированы 89 товаров, 43 РЭБ-компонента, 21 статья, 50 дополнительных материалов, "
            "12 демонстрационных схем, 4 учебных маршрута, 13 уроков и 29 практических заданий."
        ),
        "Ограничения проекта также зафиксированы:": (
            "Ограничения проекта также зафиксированы: промышленный деплой на PostgreSQL, контейнеризация, полноценная PCB-разводка, "
            "производственные форматы Gerber/Excellon, расширенная поддержка KiCad/LTspice и нейронный deep analysis относятся "
            "к следующим итерациям. Приоритет дальнейшего развития изменен на экспертный порядок: сначала объяснимые правила и факты, "
            "затем constraint/optimization и только после этого нейронные подсказки поверх проверяемого baseline."
        ),
    }
    for paragraph in doc.paragraphs:
        stripped = paragraph.text.strip()
        for prefix, text in replacements.items():
            if stripped.startswith(prefix):
                replace_paragraph(paragraph, text)

    if len(doc.tables) >= 13:
        table = doc.tables[0]
        table.cell(5, 2).text = (
            "21 статья по 6 категориям, 50 дополнительных материалов, внутренние ссылки, изображения, gif-анимации, "
            "видео, файлы, инженерная лаборатория и практикум обучения"
        )
        table.cell(6, 2).text = (
            "Схемы, CAD/SIM, проекты, версии, SimulationRun, ProjectReview, CAD-import subset, expert-first review, "
            "Pro-аналитика и rule_ai"
        )

        table = doc.tables[2]
        table.cell(4, 1).text = "POST /api/ai/chat/, self-hosted rule_ai, Expert trace из review, demo-mode без внешнего ключа"
        table.cell(4, 2).text = "Ответ строится по данным схемы, BOM, review, expert findings и статьям; внешняя LLM не является обязательной"

        table = doc.tables[5]
        rows = [
            ("Expert Review", "Design Health Score, DRC/ERC, BOM-риск, derating", "ProjectReview, expert_rules, PDF/HTML отчет"),
            ("Measurement Core", "Ожидаемое против измеренного, сохранение метрик", "ProjectMeasurement, simulation_analysis, lab_measurements"),
            ("CAD Import", "LTspice/SPICE/KiCad subset в scheme_data", "Lark parser, import preview, review после импорта"),
            ("Expert Assistant", "Объяснение ошибок и план исправления", "rule_ai, Expert trace, graph metrics, BOM facts"),
            ("Neural deep analysis", "Вероятностные подсказки после expert baseline", "PyTorch/GOLEM как отдельный будущий sprint"),
        ]
        for row_idx, values in enumerate(rows, start=1):
            for col_idx, value in enumerate(values):
                table.cell(row_idx, col_idx).text = value

        table = doc.tables[8]
        rows = [
            ("Редактор схем", "УГО-режим, GND/DRC, маршрутизация, BOM", "реализовано"),
            ("Симуляция", "DC, AC, TRAN, ngspice.wasm, JS/NumPy fallback", "реализовано"),
            ("Инженерная лаборатория", "5 расчетов и оценка норма/риск/перегрев", "реализовано"),
            ("Engineering Review", "Design Health Score, expert findings, PDF/HTML", "реализовано"),
            ("Expert-first stack", "jsonschema, rule-engine, Pint, Lark, Z3, scikit-fuzzy", "реализовано"),
            ("Практикум", "4 маршрута, 13 уроков, 29 заданий", "реализовано"),
            ("Каталог", "89 товаров, 43 РЭБ-компонента, no-Wikimedia media-policy", "реализовано"),
            ("Энциклопедия", "21 статья, 50 материалов, внутренние ссылки и связанные товары", "реализовано"),
        ]
        for row_idx, values in enumerate(rows, start=1):
            for col_idx, value in enumerate(values):
                table.cell(row_idx, col_idx).text = value

        table = doc.tables[10]
        table.cell(5, 0).text = "Наполнение и обучение"
        table.cell(5, 1).text = "Опубликованы 21 статья, 50 дополнительных материалов, 4 учебных маршрута, 13 уроков и 29 заданий"
        table.cell(5, 2).text = "На защите можно показать связку энциклопедии, лаборатории, практикума и review"

        table = doc.tables[11]
        rows = [
            ("manage.py check", "0 замечаний", "Проверка конфигурации Django"),
            ("makemigrations --check", "No changes detected", "Контроль миграций после новых сервисов"),
            ("check_demo_ready --json", "OK, включая expert_stack", "Проверка URL, learning data и библиотечных smoke-сценариев"),
            ("check_data_integrity --json", "OK, 0 errors, 0 warnings", "Аудит каталога, media, статей, ссылок и demo-схем"),
            ("Targeted regression", "18/18 expert, 16/16 learning, 8/8 search", "Проверка expert-first слоя и совместимости обучения"),
        ]
        for row_idx, values in enumerate(rows, start=1):
            for col_idx, value in enumerate(values):
                table.cell(row_idx, col_idx).text = value

        table = doc.tables[12]
        rows = [
            ("1", "Measurement Core", "Probes, сохранение измерений, expected vs measured, связь лаборатории и симулятора"),
            ("1", "Expert Review Core", "Расширение rule packs, Learning-by-review, объяснимые рекомендации"),
            ("2", "CAD Import", "LTspice/KiCad subset, import preview, автозапуск review"),
            ("2", "Комментарии пользователей", "Комментарии к товарам, урокам, статьям, demo-схемам и ProjectReview"),
            ("3", "Neural deep analysis", "PyTorch/GOLEM только после expert baseline и датасета схем"),
        ]
        for row_idx, values in enumerate(rows, start=1):
            for col_idx, value in enumerate(values):
                table.cell(row_idx, col_idx).text = value

    if not any("Приложение Е. Актуализация реализации от 19.05.2026" in p.text for p in doc.paragraphs):
        doc.add_heading("Приложение Е. Актуализация реализации от 19.05.2026", level=2)
        doc.add_paragraph(
            "После подготовки основной редакции проекта в систему добавлен expert-first контур инженерной проверки. "
            "Он не заменяет существующие CAD/SIM, лабораторию и практикум, а связывает их через общие факты, правила, "
            "единицы измерения, ограничения, импорт и оценку риска."
        )
        table = doc.add_table(rows=1, cols=3)
        for idx, value in enumerate(("Слой", "Реализация", "Использование в проекте")):
            table.rows[0].cells[idx].text = value
        rows = [
            ("Expert rules", "jsonschema + rule-engine, default_rules.json", "ProjectReview, rule_ai, demo-ready"),
            ("Unit-safe parsing", "Pint, engineering_units.py", "Лаборатория, learning grader, review"),
            ("Constraint solving", "z3-solver, constraint_solver.py", "Подбор LED, делителя, RC, NE555, стабилизатора"),
            ("CAD parsing", "Lark, cad_parsers.py", "SPICE/LTspice subset перед нормализацией в scheme_data"),
            ("Risk scoring", "scikit-fuzzy, risk_scoring.py", "Мягкая оценка риска проекта поверх экспертных фактов"),
        ]
        for values in rows:
            cells = table.add_row().cells
            for idx, value in enumerate(values):
                cells[idx].text = value
        doc.add_paragraph(
            "Порядок дальнейшего развития зафиксирован как expert systems -> constraint/optimization -> neural deep analysis. "
            "PyTorch/GOLEM, OR-Tools и RDFLib не входят в основной runtime текущей версии и отнесены к отдельным этапам развития."
        )

    saved_path = save_docx_safe(doc, diploma_path)
    Document(str(saved_path))


def speech_markdown() -> str:
    return f"""# Речь к защите и возможные вопросы по дипломной работе DOLG

Речь рассчитана примерно на 8-10 минут и следует порядку актуальной темной презентации на 12 слайдов. Подробности оставлены в речи, а на слайдах информация разложена по панелям, чтобы не перегружать экран.

## Речь по слайдам

### Слайд 1. Титульный слайд
Добрый день. Тема моей дипломной работы — «{TITLE}». В работе рассматривается веб-сервис, который связывает каталог компонентов, инженерные расчеты, построение схем, симуляцию, обучение и оформление заказа.

Главная идея проекта состоит в том, чтобы пользователь не переносил данные между разрозненными системами вручную. В DOLG компонент выбирается в каталоге, используется в схеме, проверяется расчетом и симуляцией, а затем попадает в BOM и заказ. Поэтому на первом слайде оставлены тема, авторские данные, актуальный экран каталога и короткая логика проекта.

### Слайд 2. Актуальность и цель
Актуальность связана с тем, что реальные инженерные действия часто разбиты между несколькими средами: каталогом, datasheet, CAD, SPICE-симулятором, таблицей BOM и системой заказа. При таком подходе теряется связь между номиналом, моделью, документацией и расчетным результатом.

Цель работы — разработать веб-приложение DOLG для подбора и приобретения радио- и электронных компонентов, проектирования схем, инженерных расчетов, симуляции, обучения и проверки результата в едином интерфейсе. В текущем состоянии проект содержит 89 товаров, 5 инженерных расчетов, 4 учебных маршрута, 13 уроков, 29 заданий и 12 демонстрационных схем.

### Слайд 3. Анализ решений
На слайде показаны классы аналогов. Интернет-магазины хорошо решают поиск и заказ, но не проверяют схему. Онлайн-симуляторы помогают быстро экспериментировать, но обычно работают с абстрактными элементами. Профессиональные CAD/EDA-системы, такие как KiCad или Altium, дают мощный контур проектирования, но имеют высокий порог входа и не связаны напрямую с учебным маршрутом и покупкой компонентов.

DOLG занимает промежуточную нишу: это не замена промышленной PCB CAD-системы, а учебно-инженерный веб-сервис. Его ценность в том, что компонент, схема, расчет, измерение, обучение и заказ оказываются связаны одним маршрутом.

### Слайд 4. Целевая аудитория и сценарии
Проект ориентирован на студентов, радиолюбителей и инженеров. Студентам важна быстрая обратная связь: расчет, схема, автопроверка и объяснение ошибки. Радиолюбителям нужен быстрый путь от идеи до набора компонентов. Инженерам полезны оценка номиналов, теплового запаса, измерений и повторное использование проверенных решений.

Типовой маршрут выглядит так: пользователь находит компонент, открывает карточку и документацию, выполняет расчет в лаборатории, собирает схему, запускает симуляцию, сравнивает ожидаемое и измеренное, затем формирует BOM и заказ. Такой сценарий показывает, что сайт является не только витриной товаров, а рабочим контуром.

### Слайд 5. Архитектура системы
Архитектура построена на Django. Логика разделена между приложениями shop, accounts, orders, knowledge и Dolg_APP. При этом расчеты, проверки, review, импорт и обучение вынесены в общий service-layer. Это важно, потому что формулы и проверки не дублируются в шаблонах или JavaScript.

В текущей версии особенно важен expert-first слой. Он включает rule packs, проверку правил через rule-engine, единицы измерения через Pint, разбор SPICE/LTspice subset через Lark, подбор номиналов через Z3 и мягкую оценку риска через scikit-fuzzy. Нейронная сеть в эту часть специально не встроена: сначала должны работать объяснимые инженерные правила.

### Слайд 6. Ключевые фрагменты кода
На слайде показаны три опорные структуры. Первая — экспертный finding: у него есть rule_id, severity, evidence, recommendation и confidence. Это позволяет объяснить, почему система считает схему рискованной или неполной.

Вторая структура — unit-safe parser. Он приводит значения вроде 6.8kOhm, 2.5 мА или 100 нФ к единому численному виду. Это защищает лабораторию, review и учебные задания от расхождений в единицах измерения. Третья структура — constraint solver на Z3. Он не заменяет инженера, а подбирает допустимые варианты номиналов при заданных ограничениях.

### Слайд 7. Реализованные модули
На этом слайде функции отделены от свойств. Ассортимент и единые карточки — это свойства каталога, а функциональные модули строятся поверх них. Реализованы каталог и BOM, инженерная лаборатория, практикум обучения, CAD/SIM, Engineering Review и контроль качества.

Особенно важно, что новые части используют общие сервисы. Например, расчет из лаборатории может стать эталоном в учебной задаче, а проверка схемы может попасть в Engineering Review и в ответ self-hosted AI-помощника.

### Слайд 8. CAD-редактор
CAD-редактор выделен отдельным слайдом, потому что это отдельный режим работы. Он основан на HTML5 Canvas2D и поддерживает сетку, snap-привязку, слои, элементы чертежа и ГОСТ-шаблоны.

На слайде оставлен визуальный акцент на рабочей области, а подробности перенесены в речь. Редактор нужен не ради декоративного рисования, а ради подготовки схемы и технических материалов, которые затем можно передать в симулятор, review и демо-сценарии.

### Слайд 9. Симуляция и измерения
Симулятор поддерживает DC, AC и TRAN-анализ. Основной расчетный контур использует ngspice.wasm, а для простых DC-цепей предусмотрен server-side fallback. Дополнительно добавлен Pro-слой аналитики: FFT-спектр, Bode plot, Monte Carlo tolerance, THD/SINAD/ENOB, parameter sweep.

Важное развитие — измерительный контур. Пользователь может работать не только с графиком, но и с метриками: напряжением узла, током ветви, RMS, частотой, duty cycle, мощностью и температурой. Эти же метрики используются в лаборатории и практикуме.

### Слайд 10. Проверка результата
Проверка результата строится на нескольких уровнях. manage.py check проходит без ошибок, миграции не требуют новых изменений, demo-ready и data-integrity возвращают OK. Отдельно проверяются no-Wikimedia media-policy, опубликованные учебные данные, demo-схемы и expert stack.

На слайде справа показан результат, который пользователь получает в инженерной лаборатории: расчет не ограничивается числом, а сопровождается оценкой «норма» или «риск». Это важно для защиты, потому что демонстрирует прикладной инженерный результат, а не только внутренние тесты.

### Слайд 11. План развития
План развития теперь строится по принципу: крупная фича, обучающий блок и документация. Ближайшие направления — усиление Engineering Review, развитие Measurement Core, расширение CAD Import, добавление комментариев зарегистрированных пользователей и только затем neural deep analysis.

Нейронная сеть не должна заменять экспертную систему. Сначала проект должен объяснять вывод через правила, факты, формулы, граф схемы, измерения и BOM. Нейронный слой можно подключать позже как вероятностную подсказку deep_hint, обязательно сравнивая ее с экспертным baseline.

### Слайд 12. Спасибо за внимание
В результате работы создана демонстрационная веб-платформа, которая объединяет каталог компонентов, CAD/SIM, инженерную лабораторию, практикум и экспертную проверку проекта. Главное отличие DOLG — не отдельная функция, а связанный маршрут: компонент → расчет → схема → измерение → review → BOM → заказ.

Спасибо за внимание. Готов ответить на вопросы.

## Возможные вопросы и ответы

1. **Почему выбран порядок expert-first, а не сразу нейронная сеть?**
   Потому что инженерное решение должно быть объяснимым. Сначала система должна показывать правило, факты, расчет и рекомендацию. Нейронная сеть может появиться позже как слой глубокого анализа, но не как источник финального verdict.

2. **Что дает rule-engine в проекте?**
   Он позволяет хранить условия проверок в rule packs и получать выводы с rule_id, severity, evidence и recommendation. Это удобнее, чем разбрасывать условия по view-логике.

3. **Зачем нужен Pint?**
   Пользователь может вводить значения в разных формах: 10k, 6.8kOhm, 2.5мА, 100нФ. Pint и общий unit-service приводят их к единому виду для лаборатории, review и обучения.

4. **Зачем нужен Z3, если расчеты можно сделать формулой?**
   Формула дает одно значение, а constraint solver подбирает допустимые варианты при ограничениях: например, номиналы делителя или LED-резистора с учетом диапазонов.

5. **Чем DOLG отличается от KiCad или Altium?**
   DOLG не заменяет промышленный CAD. Он закрывает учебно-практический веб-маршрут, где компонент связан с расчетом, симуляцией, обучением и заказом.

6. **Почему в демонстрации важен no-Wikimedia media-policy?**
   Потому что карточки товаров должны выглядеть контролируемо и не получать случайные изображения. Активные карточки используют локальные assets или сгенерированные технические заглушки.

7. **Как проверяется готовность проекта к показу?**
   Используются manage.py check, makemigrations --check --dry-run, check_demo_ready --json, check_data_integrity --json и targeted-тесты по экспертному слою, обучению и поиску.

8. **Что такое Engineering Review?**
   Это единый отчет по проекту: схема, DRC/ERC, BOM-риск, наличие GND/источника, derating, измерения, экспертные findings и итоговая оценка готовности.

9. **Почему обучение вынесено в knowledge?**
   Потому что обучение связано с энциклопедией, лабораторией и статьями. При этом старый /learn/ сохранен как redirect на /knowledge/learning/.

10. **Какие типы учебных заданий поддерживаются?**
   math_numeric для числового ответа, circuit_build для сборки схемы и simulation_measure для проверки измеренного результата симуляции.

11. **Что осталось развивать в первую очередь?**
   Связать лабораторию и симулятор еще плотнее: probes, сохраненные измерения, expected vs measured и автоматическую отправку результата в учебные задания.

12. **Когда стоит подключать PyTorch/GOLEM?**
   После стабилизации expert core и накопления датасета схем. Нейронка должна сравниваться с rule-based baseline и иметь fallback.

13. **Почему комментарии важны для будущего?**
   Зарегистрированные пользователи смогут обсуждать товары, статьи, уроки, demo-схемы и review. Это превратит сайт из витрины в учебно-инженерное пространство.

14. **Можно ли использовать проект в продакшне прямо сейчас?**
   Проект готов к локальной демонстрации и защите. Для продакшна нужны PostgreSQL, деплой, мониторинг, резервные копии и регламент работы с media/static.

15. **Главный результат работы?**
   Создана платформа, где электронная коммерция и инженерная проверка не разорваны, а связаны одним пользовательским маршрутом.
"""


def update_speech() -> None:
    md = speech_markdown()
    md_path = first_file("Речь_и_вопросы_к_защите_DOLG_*.md")
    docx_path = first_file("Речь_и_вопросы_к_защите_DOLG_*.docx")
    md_path.write_text(md, encoding="utf-8")

    doc = Document()
    doc.styles["Normal"].font.name = "Times New Roman"
    doc.styles["Normal"].font.size = DocPt(12)
    for line in md.splitlines():
        if not line.strip():
            continue
        if line.startswith("# "):
            paragraph = doc.add_heading(line[2:].strip(), level=1)
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        elif line.startswith("## "):
            doc.add_heading(line[3:].strip(), level=2)
        elif line.startswith("### "):
            doc.add_heading(line[4:].strip(), level=3)
        else:
            paragraph = doc.add_paragraph()
            add_run(paragraph, line.strip())
    saved_path = save_docx_safe(doc, docx_path)
    Document(str(saved_path))


# ---------------- PPTX ----------------


class DarkDeck:
    bg = RGBColor(7, 18, 40)
    panel = RGBColor(13, 31, 64)
    panel2 = RGBColor(17, 41, 78)
    stroke = RGBColor(0, 188, 214)
    title = RGBColor(236, 248, 255)
    text = RGBColor(197, 215, 237)
    muted = RGBColor(132, 157, 190)
    cyan = RGBColor(0, 211, 232)
    blue = RGBColor(64, 139, 255)
    green = RGBColor(32, 204, 132)
    orange = RGBColor(255, 117, 58)
    purple = RGBColor(158, 110, 255)
    red = RGBColor(255, 78, 90)

    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
        self.w = self.prs.slide_width
        self.h = self.prs.slide_height

    @staticmethod
    def e(value: float):
        return Inches(value)

    def slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.bg
        top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.w, self.e(0.08))
        top.fill.solid()
        top.fill.fore_color.rgb = self.cyan
        top.line.fill.background()
        return slide

    def text_box(self, slide, text, x, y, w, h, *, size=12, color=None, bold=False, align="left", font="Segoe UI"):
        shape = slide.shapes.add_textbox(self.e(x), self.e(y), self.e(w), self.e(h))
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = self.e(0.06)
        tf.margin_right = self.e(0.06)
        tf.margin_top = self.e(0.03)
        tf.margin_bottom = self.e(0.03)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        color = color or self.text
        for idx, line in enumerate(str(text).split("\n")):
            paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            paragraph.text = line
            paragraph.alignment = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
            for run in paragraph.runs:
                run.font.name = font
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = color
        return shape

    def title_block(self, slide, title, subtitle="", num=None):
        self.text_box(slide, title, 0.72, 0.38, 9.7, 0.5, size=25, color=self.title, bold=True)
        if subtitle:
            self.text_box(slide, subtitle, 0.74, 0.9, 12.8, 0.35, size=11.5, color=self.muted)
        if num:
            self.text_box(slide, f"{num}/12", 14.85, 8.48, 0.45, 0.2, size=8, color=self.muted, align="right")
            self.text_box(slide, "DOLG · дипломная работа", 0.75, 8.48, 2.4, 0.2, size=8, color=self.muted)

    def card(self, slide, x, y, w, h, title="", body="", accent=None, *, title_size=13, body_size=9.7, fill=None):
        accent = accent or self.cyan
        fill = fill or self.panel
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, self.e(x), self.e(y), self.e(w), self.e(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = RGBColor(25, 83, 120)
        shape.line.width = Pt(0.9)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.e(x), self.e(y), self.e(0.055), self.e(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        if title:
            self.text_box(slide, title, x + 0.18, y + 0.16, w - 0.3, 0.34, size=title_size, color=self.title, bold=True)
        if body:
            self.text_box(slide, body, x + 0.18, y + 0.58, w - 0.3, h - 0.68, size=body_size, color=self.text)
        return shape

    def metric(self, slide, x, y, w, h, value, label, color=None):
        color = color or self.cyan
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, self.e(x), self.e(y), self.e(w), self.e(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.panel2
        shape.line.color.rgb = color
        shape.line.width = Pt(1)
        self.text_box(slide, value, x, y + 0.08, w, 0.32, size=18, color=color, bold=True, align="center")
        self.text_box(slide, label, x + 0.04, y + 0.48, w - 0.08, 0.26, size=7.8, color=self.muted, align="center")

    def image(self, slide, path: Path, x, y, w, h, label=""):
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, self.e(x), self.e(y), self.e(w), self.e(h))
        panel.fill.solid()
        panel.fill.fore_color.rgb = self.panel
        panel.line.color.rgb = RGBColor(24, 93, 129)
        panel.line.width = Pt(1)
        label_h = 0.32 if label else 0
        if label:
            self.text_box(slide, label, x + 0.2, y + 0.12, w - 0.4, 0.22, size=8.5, color=self.muted)
        with Image.open(path) as image:
            iw, ih = image.size
        max_w = w - 0.32
        max_h = h - 0.32 - label_h
        scale = min(max_w / iw, max_h / ih)
        pw = iw * scale
        ph = ih * scale
        px = x + (w - pw) / 2
        py = y + label_h + (h - label_h - ph) / 2
        slide.shapes.add_picture(str(path), self.e(px), self.e(py), width=self.e(pw), height=self.e(ph))

    def table(self, slide, x, y, w, h, headers, rows, widths=None):
        shape = slide.shapes.add_table(len(rows) + 1, len(headers), self.e(x), self.e(y), self.e(w), self.e(h))
        table = shape.table
        if widths:
            for idx, width in enumerate(widths):
                table.columns[idx].width = self.e(width)
        for c, header in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(15, 60, 92)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Segoe UI"
                    run.font.size = Pt(8.2)
                    run.font.bold = True
                    run.font.color.rgb = self.title
        for r, row in enumerate(rows, 1):
            for c, value in enumerate(row):
                cell = table.cell(r, c)
                cell.text = value
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.panel
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "Segoe UI"
                        run.font.size = Pt(7.2)
                        run.font.color.rgb = self.text


def build_presentation() -> None:
    deck = DarkDeck()
    img_base = DOCS / "diploma_assets" / "screenshots" / "presentation_v6"
    catalog = img_base / "crops" / "catalog_cards_crop.jpg"
    lab = img_base / "crops" / "engineering_lab_crop.jpg"
    cad = img_base / "crops" / "cad_editor_crop.jpg"
    sim = img_base / "crops" / "simulation_ac_crop.jpg"

    # 1
    s = deck.slide()
    deck.text_box(s, "Дипломная работа", 0.78, 0.72, 5.6, 0.62, size=31, color=deck.title, bold=True)
    deck.text_box(s, TITLE, 0.82, 1.55, 5.95, 1.45, size=18, color=deck.title, bold=True)
    deck.text_box(s, f"Автор: {AUTHOR}\nНаучный руководитель: {SUPERVISOR}\nГод защиты: 2026", 0.86, 6.55, 5.5, 0.82, size=11.8, color=deck.text)
    deck.image(s, catalog, 7.05, 0.68, 8.15, 6.9, "Актуальный каталог с карточками компонентов")
    deck.text_box(s, "каталог → расчёт → схема → измерение → review → BOM → заказ", 7.35, 7.78, 7.55, 0.28, size=10.7, color=deck.cyan, align="center", bold=True)

    # 2
    s = deck.slide(); deck.title_block(s, "Актуальность и цель", "Компонент, схема, расчёт и заказ должны жить в одном инженерном маршруте", 2)
    deck.card(s, 0.75, 1.55, 4.55, 2.35, "Разрыв инструментов", "Каталог, datasheet, CAD, SPICE, BOM и заказ часто работают отдельно. Из-за ручного переноса теряются номиналы, модели, документация и результаты измерений.", deck.blue, body_size=10.6)
    deck.card(s, 5.7, 1.55, 4.55, 2.35, "Практическая потребность", "Учебным и малым инженерным проектам нужен быстрый web-путь: подобрать компонент, проверить схему, получить оценку риска и сформировать заказ.", deck.orange, body_size=10.6)
    deck.card(s, 10.65, 1.55, 4.55, 2.35, "Цель работы", "Разработать DOLG для подбора, покупки, проектирования, симуляции, обучения и инженерной проверки результата в едином интерфейсе.", deck.green, body_size=10.6)
    for idx, data in enumerate([("89", "товаров"), ("43", "РЭБ"), ("5", "расчётов"), ("4", "маршрута"), ("13", "уроков"), ("29", "заданий"), ("12", "demo")]):
        deck.metric(s, 0.9 + idx * 2.08, 5.05, 1.55, 0.82, data[0], data[1], [deck.blue, deck.cyan, deck.green, deck.purple, deck.orange, deck.red, deck.cyan][idx])
    deck.card(s, 1.05, 6.65, 13.9, 0.95, "Формулировка результата", "DOLG показывает связку электронной коммерции и инженерной проверки: пользователь выбирает реальный компонент, применяет его в схеме, проверяет расчетом/симуляцией и получает BOM для заказа.", deck.cyan, title_size=12, body_size=10.2)

    # 3
    s = deck.slide(); deck.title_block(s, "Анализ решений", "Позиционирование относительно магазинов, онлайн-симуляторов и CAD/EDA", 3)
    deck.table(s, 0.75, 1.48, 14.45, 4.25, ["Класс", "Сильная сторона", "Ограничение", "DOLG"], [
        ("Интернет-магазины", "каталог, цена, наличие, заказ", "нет проверки схемы", "BOM и товарная карточка связаны со схемой"),
        ("Онлайн-симуляторы", "быстрые учебные эксперименты", "абстрактные элементы", "схема работает рядом с реальными позициями"),
        ("CAD / EDA", "полный цикл разработки", "порог входа и отдельная закупка", "упрощенный web-маршрут для учебных задач"),
        ("AI-CAD сервисы", "подсказки поверх схемы", "риск непрозрачных выводов", "expert-first trace перед нейронным анализом"),
        ("DOLG", "каталог + CAD/SIM + lab + learning + review", "демонстрационный масштаб ВКР", "единый маршрут от компонента до заказа"),
    ], widths=[2.6, 3.7, 3.4, 4.75])
    deck.card(s, 1.0, 6.25, 13.95, 1.15, "Практическая ниша", "Проект не конкурирует с промышленными PCB CAD напрямую. Он закрывает сценарий, который редко собран в одном месте: компонент → расчёт → схема → измерение → обучение → review → BOM → заказ.", deck.cyan, title_size=12.5, body_size=10.8)

    # 4
    s = deck.slide(); deck.title_block(s, "Целевая аудитория и сценарии", "Три группы пользователей и единый рабочий маршрут", 4)
    deck.card(s, 0.75, 1.45, 4.25, 1.75, "Студенты", "Изучают схемотехнику через расчёты, готовые схемы, практикум и автопроверку. Важны подсказки и обратная связь.", deck.blue, body_size=10)
    deck.card(s, 0.75, 3.45, 4.25, 1.75, "Радиолюбители", "Подбирают компоненты под идею, быстро проверяют схему и собирают набор деталей для покупки.", deck.green, body_size=10)
    deck.card(s, 0.75, 5.45, 4.25, 1.75, "Инженеры", "Сравнивают номиналы, оценивают запас, сохраняют проверенные схемы и экспортируют BOM.", deck.orange, body_size=10)
    deck.card(s, 5.55, 1.45, 9.65, 3.1, "Типовой маршрут", "1. Найти компонент → 2. Открыть карточку и datasheet → 3. Рассчитать узел в лаборатории → 4. Собрать схему → 5. Запустить DC/AC/TRAN → 6. Сравнить ожидаемое и измеренное → 7. Получить review → 8. Сформировать BOM и заказ.", deck.cyan, title_size=15, body_size=13)
    for idx, data in enumerate([("21", "статья"), ("50", "материалов"), ("4", "tracks"), ("13", "lessons"), ("29", "tasks")]):
        deck.metric(s, 5.75 + idx * 1.82, 5.45, 1.4, 0.85, data[0], data[1], [deck.blue, deck.green, deck.purple, deck.orange, deck.red][idx])
    deck.text_box(s, "Маршрут на слайде показывает, что обучение, лаборатория и review не отдельные витрины, а продолжение работы с компонентом.", 5.8, 6.75, 8.75, 0.48, size=10, color=deck.muted)

    # 5
    s = deck.slide(); deck.title_block(s, "Архитектура системы", "Django apps + service-layer + expert-first контур", 5)
    deck.card(s, 0.65, 1.35, 3.45, 4.85, "Django apps", "shop — каталог, поиск, BOM\naccounts — профиль и доступ\norders — корзина и заказы\nknowledge — статьи, lab, learning\nDolg_APP — CAD/SIM, проекты, review", deck.blue)
    deck.card(s, 4.25, 1.35, 3.6, 4.85, "Service-layer", "engineering_lab\nlearning_grader\nsimulation_analysis\nschematic_graph\nproject_review\ncad_import\nproduct_images", deck.green)
    deck.card(s, 8.05, 1.35, 3.6, 4.85, "Expert-first core", "jsonschema rule packs\nrule-engine predicates\nPint units\nLark parsers\nZ3 constraints\nscikit-fuzzy risk\nExpert trace for AI", deck.orange)
    deck.card(s, 11.85, 1.35, 3.45, 4.85, "Client / UI", "каталог и карточки\nCanvas2D CAD/SIM\nngspice.wasm + fallback\nPro-графики\nлаборатория\nпрактикум и отчёты", deck.purple)
    deck.card(s, 1.05, 6.65, 13.95, 0.8, "Принцип развития", "Крупная фича добавляется вместе с обучающим блоком и документацией; формулы, проверки и метрики остаются в Python service-layer.", deck.cyan, title_size=12, body_size=10)

    # 6
    s = deck.slide(); deck.title_block(s, "Ключевые фрагменты кода", "Кодовые структуры, через которые объясняется инженерная логика проекта", 6)
    code1 = "finding = {\n  'rule_id': rule['id'],\n  'severity': rule['severity'],\n  'evidence': evidence,\n  'recommendation': text,\n  'confidence': confidence,\n}"
    code2 = "parse_engineering_quantity(\n  '6.8kOhm', expected_unit='Ом'\n) -> 6800\n\nparse_engineering_quantity(\n  '2.5мА', expected_unit='А'\n) -> 0.0025"
    code3 = "solve_design_constraints(\n  'voltage_divider',\n  {'vin': 9, 'target_vout': 3}\n)\n# допустимые пары R1/R2\n# и объяснение ограничений"
    deck.card(s, 0.72, 1.45, 4.7, 5.75, "Expert finding", "", deck.blue)
    deck.text_box(s, code1, 1.0, 2.05, 4.1, 2.25, size=12, color=deck.title, font="Consolas")
    deck.text_box(s, "Review и AI получают не просто строку ошибки, а проверяемый вывод: правило, факты, уверенность и рекомендацию.", 1.0, 5.35, 4.0, 0.82, size=10, color=deck.text)
    deck.card(s, 5.65, 1.45, 4.7, 5.75, "Unit-safe layer", "", deck.green)
    deck.text_box(s, code2, 5.92, 2.05, 4.1, 2.25, size=11.5, color=deck.title, font="Consolas")
    deck.text_box(s, "Pint снимает конфликт единиц между лабораторией, review и учебными задачами.", 5.92, 5.35, 4.0, 0.82, size=10, color=deck.text)
    deck.card(s, 10.58, 1.45, 4.7, 5.75, "Constraint solver", "", deck.orange)
    deck.text_box(s, code3, 10.85, 2.05, 4.1, 2.45, size=11, color=deck.title, font="Consolas")
    deck.text_box(s, "Z3 подбирает варианты при ограничениях; финальный инженерный verdict остаётся у rule/review слоя.", 10.85, 5.35, 4.0, 0.82, size=10, color=deck.text)

    # 7
    s = deck.slide(); deck.title_block(s, "Реализованные модули", "Функции проекта поверх общего каталога, данных и сервисного слоя", 7)
    modules = [
        ("Каталог и BOM", "89 товаров, 43 РЭБ, фильтры, datasheet, сравнение, XLSX/CSV, add-all-to-cart.", deck.blue),
        ("Инженерная лаборатория", "Транзисторный ключ, NE555, стабилизатор, RC-антидребезг, тепловой запас.", deck.green),
        ("Практикум", "4 маршрута, 13 уроков, 29 заданий: math_numeric, circuit_build, simulation_measure.", deck.purple),
        ("CAD/SIM", "Canvas2D редактор, УГО, DRC, ngspice.wasm, DC/AC/TRAN и Pro-аналитика.", deck.orange),
        ("Engineering Review", "Design Health Score, DRC/ERC, BOM-риск, derating, expert findings, PDF/HTML.", deck.cyan),
        ("Контроль качества", "check, migrations, demo-ready, data-integrity, no-Wikimedia и targeted tests.", deck.red),
    ]
    for i, (title, body, color) in enumerate(modules):
        deck.card(s, 0.75 + (i % 3) * 5.0, 1.42 + (i // 3) * 2.45, 4.55, 1.95, title, body, color, title_size=13, body_size=9.6)
    deck.card(s, 1.0, 6.72, 13.9, 0.76, "Важно", "Ассортимент и единые карточки являются свойствами каталога, а функции — это инженерные сценарии поверх данных: lab, CAD/SIM, learning, review, BOM и order flow.", deck.cyan, title_size=12, body_size=9.8)

    # 8
    s = deck.slide(); deck.title_block(s, "CAD-редактор", "Отдельный режим для чертежей, УГО, ГОСТ-шаблонов и передачи схемы в симулятор", 8)
    deck.image(s, cad, 0.75, 1.35, 9.5, 6.55, "Рабочая область CAD: сетка, ГОСТ-шаблон, примитивы и слои")
    deck.card(s, 10.55, 1.42, 4.65, 1.35, "Режим CAD", "Canvas2D, сетка, snap-привязка, слои, аккуратное позиционирование.", deck.blue)
    deck.card(s, 10.55, 3.02, 4.65, 1.35, "ГОСТ и чертежи", "Рамки, титульные блоки, размерные линии и элементы ЕСКД-логики.", deck.green)
    deck.card(s, 10.55, 4.62, 4.65, 1.35, "Библиотека УГО", "Символы компонентов, узлы, подписи и подготовка схемы к расчету.", deck.orange)
    deck.card(s, 10.55, 6.22, 4.65, 1.35, "Экспорт", "PDF, PNG, SVG, JSON; JSON идет в симуляцию, review и demo-сценарии.", deck.purple)

    # 9
    s = deck.slide(); deck.title_block(s, "Симуляция и измерения", "Схема превращается в расчёт, графики, метрики и инженерную оценку", 9)
    deck.card(s, 0.75, 1.35, 4.65, 1.35, "Режимы анализа", "DC / OP — рабочая точка\nAC — АЧХ/ФЧХ и -3 дБ\nTRAN — переходные процессы", deck.blue, body_size=9.2)
    deck.card(s, 0.75, 2.98, 4.65, 1.55, "Pro-аналитика", "FFT spectrum, Bode plot, Monte Carlo tolerance, THD/SINAD/ENOB, parameter sweep, server-side fallback.", deck.green, body_size=9.4)
    deck.card(s, 0.75, 4.82, 4.65, 1.62, "Метрики", "Напряжение узла, ток ветви, RMS, частота, duty cycle, мощность элемента и температура.", deck.orange, body_size=9.6)
    deck.card(s, 0.75, 6.73, 4.65, 0.78, "Связь с обучением", "Expected vs measured используется в лаборатории, review и simulation_measure задачах.", deck.cyan, title_size=11.4, body_size=8.8)
    deck.image(s, sim, 5.72, 1.35, 9.55, 6.55, "Результаты симуляции: графики, измерения и быстрые действия")

    # 10
    s = deck.slide(); deck.title_block(s, "Проверка результата", "Тесты, demo-ready, data-integrity и лабораторный результат", 10)
    deck.table(s, 0.75, 1.35, 5.35, 3.05, ["Уровень", "Результат"], [
        ("manage.py check", "0 issues"),
        ("makemigrations", "No changes detected"),
        ("check_demo_ready", "OK, включая expert_stack"),
        ("check_data_integrity", "OK, 0 errors / 0 warnings"),
        ("Targeted tests", "18/18 expert, 16/16 learning, 8/8 search"),
    ], widths=[2.15, 3.2])
    for i, (v, l, c) in enumerate([("OK", "demo-ready", deck.green), ("OK", "integrity", deck.green), ("0", "warnings", deck.blue), ("89", "products", deck.cyan)]):
        deck.metric(s, 0.85 + i * 1.3, 4.85, 1.02, 0.75, v, l, c)
    deck.card(s, 0.8, 6.08, 5.2, 1.2, "Что подтверждается", "Каталог, media-policy, статьи, demo-схемы, learning, expert stack и лабораторные расчеты проверяются отдельными smoke/regression сценариями.", deck.cyan, title_size=11.8, body_size=9.3)
    deck.image(s, lab, 6.45, 1.28, 8.8, 6.65, "Инженерная лаборатория: расчет + оценка «норма/риск»")

    # 11
    s = deck.slide(); deck.title_block(s, "План развития", "Пакетный подход: фича + обучалка + документация", 11)
    road = [
        ("1. Expert Review Core", "расширение rule packs, explainable findings, Learning-by-review", deck.blue),
        ("2. Measurement Core", "probes, сохранение измерений, expected vs measured, sweep", deck.green),
        ("3. CAD Import", "LTspice/KiCad subset, import preview, review после импорта", deck.orange),
        ("4. Комментарии", "товары, уроки, статьи, demo-схемы, ProjectReview", deck.purple),
        ("5. Neural deep analysis", "PyTorch/GOLEM позже: deep_hint поверх expert baseline", deck.cyan),
    ]
    for i, (title, body, color) in enumerate(road):
        deck.card(s, 0.85, 1.35 + i * 1.15, 6.05, 0.9, title, body, color, title_size=11.6, body_size=8.8)
    deck.card(s, 7.45, 1.45, 7.25, 5.15, "Главный принцип следующего этапа", "1. Экспертная система объясняет вывод через правила, факты, формулы, граф схемы, измерения и BOM.\n\n2. Constraint/optimization подбирает допустимые варианты номиналов и BOM.\n\n3. Нейронный слой появляется только как поздний анализ и не выносит финальный инженерный verdict.", deck.cyan, title_size=15, body_size=12.7)
    deck.text_box(s, "Такой порядок снижает риск непрозрачного AI и делает выводы пригодными для защиты.", 7.7, 6.95, 6.7, 0.4, size=10.5, color=deck.muted, align="center")

    # 12
    s = deck.slide()
    deck.text_box(s, "Спасибо за внимание!", 0.85, 0.92, 8.5, 0.75, size=33, color=deck.title, bold=True)
    deck.text_box(s, "DOLG объединяет каталог, CAD/SIM, инженерную лабораторию, практикум и экспертную проверку проекта в одном тёмном рабочем интерфейсе.", 0.9, 2.0, 7.4, 0.95, size=16, color=deck.text)
    deck.card(s, 0.95, 4.15, 5.6, 2.05, "Контактная информация", f"{AUTHOR}\nEmail: buryako@internet.com\nТелефон: +7 (903) 439-44-87", deck.cyan, title_size=14, body_size=12.5)
    deck.card(s, 7.25, 4.15, 7.05, 2.05, "К защите готово", "Диплом, речь и презентация актуализированы. Demo-ready и data-integrity проходят; expert-first stack зафиксирован в документах.", deck.green, title_size=14, body_size=12.2)
    deck.text_box(s, "Вопросы?", 5.75, 7.22, 4.5, 0.55, size=25, color=deck.cyan, bold=True, align="center")
    deck.title_block(s, "", "", 12)

    docs_pptx = first_file("Презентация_DOLG_*.pptx")
    actual_pptx = DOCS / "Презентация_DOLG_актуальная_20260519.pptx"
    saved_docs_pptx = save_pptx_safe(deck.prs, docs_pptx)
    saved_actual_pptx = save_pptx_safe(deck.prs, actual_pptx)
    download_target = DOWNLOADS / "Razrabotka-veb-prilozheniya-dlya-prodazhi-radio-i-elektronnyh-komponentov-so-vstroennymi-instrumenta_updated_20260518.pptx"
    download_fixed = DOWNLOADS / "Razrabotka-veb-prilozheniya-dlya-prodazhi-radio-i-elektronnyh-komponentov-so-vstroennymi-instrumenta_fixed_dark_20260519.pptx"
    if download_target.exists():
        copy_file_safe(saved_docs_pptx, download_target)
    copy_file_safe(saved_docs_pptx, download_fixed)
    Presentation(str(saved_docs_pptx))
    Presentation(str(saved_actual_pptx))


def main() -> None:
    backup_dir = DOCS / "presentation_backups" / "20260519_dark_fix_before_script"
    backup_dir.mkdir(parents=True, exist_ok=True)
    for path in [
        first_file("Диплом*.docx"),
        first_file("Речь_и_вопросы_к_защите_DOLG_*.docx"),
        first_file("Речь_и_вопросы_к_защите_DOLG_*.md"),
        first_file("Презентация_DOLG_*.pptx"),
    ]:
        backup_file(path, backup_dir)
    download_pptx = DOWNLOADS / "Razrabotka-veb-prilozheniya-dlya-prodazhi-radio-i-elektronnyh-komponentov-so-vstroennymi-instrumenta_updated_20260518.pptx"
    backup_file(download_pptx, backup_dir)

    update_diploma()
    update_speech()
    build_presentation()
    print("Updated diploma, speech and dark presentation.")


if __name__ == "__main__":
    main()
