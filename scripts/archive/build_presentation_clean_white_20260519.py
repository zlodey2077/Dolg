from __future__ import annotations

import shutil
from pathlib import Path

from docx import Document
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / 'docs'
DOWNLOADS = Path.home() / 'Downloads'
ASSETS = DOCS / 'diploma_assets' / 'screenshots' / 'presentation_v6' / 'crops'

TITLE = (
    'Разработка веб-приложения для продажи радио- и электронных компонентов '
    'со встроенными инструментами проектирования и симуляции схем'
)
AUTHOR = 'Буряко Дмитрий Сергеевич'
SUPERVISOR = 'Буланов Сергей Георгиевич'


def first_file(pattern: str) -> Path:
    files = [path for path in DOCS.glob(pattern) if not path.name.startswith('~$')]
    if not files:
        raise FileNotFoundError(pattern)
    return files[0]


def backup_file(path: Path, backup_dir: Path) -> None:
    if path.exists():
        target = backup_dir / path.name
        if not target.exists():
            shutil.copy2(path, target)


def copy_file_safe(source: Path, target: Path) -> Path:
    try:
        shutil.copy2(source, target)
        return target
    except PermissionError:
        fallback = target.with_name(f'{target.stem}_white_20260519{target.suffix}')
        shutil.copy2(source, fallback)
        return fallback


def update_speech_intro() -> None:
    md_path = first_file('Речь_и_вопросы_к_защите_DOLG_*.md')
    text = md_path.read_text(encoding='utf-8')
    text = text.replace('актуальной темной презентации', 'актуальной светлой презентации')
    text = text.replace('актуальной тёмной презентации', 'актуальной светлой презентации')
    md_path.write_text(text, encoding='utf-8')

    docx_candidates = sorted(
        [p for p in DOCS.glob('Речь_и_вопросы_к_защите_DOLG_*.docx') if not p.name.startswith('~$')],
        key=lambda p: p.stat().st_mtime,
        reverse=True,
    )
    if not docx_candidates:
        return
    docx_path = docx_candidates[0]
    doc = Document(str(docx_path))
    for paragraph in doc.paragraphs:
        if (
            'актуальной темной презентации' in paragraph.text
            or 'актуальной тёмной презентации' in paragraph.text
        ):
            paragraph.text = paragraph.text.replace(
                'актуальной темной презентации', 'актуальной светлой презентации'
            )
            paragraph.text = paragraph.text.replace(
                'актуальной тёмной презентации', 'актуальной светлой презентации'
            )
    try:
        doc.save(str(docx_path))
    except PermissionError:
        fallback = docx_path.with_name(f'{docx_path.stem}_white_20260519{docx_path.suffix}')
        doc.save(str(fallback))


class CleanDeck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
        self.W = 16
        self.H = 9
        self.bg = RGBColor(248, 250, 252)
        self.white = RGBColor(255, 255, 255)
        self.navy = RGBColor(15, 32, 59)
        self.text = RGBColor(39, 54, 77)
        self.muted = RGBColor(96, 111, 132)
        self.line = RGBColor(214, 224, 235)
        self.cyan = RGBColor(0, 151, 190)
        self.blue = RGBColor(37, 99, 235)
        self.green = RGBColor(22, 163, 74)
        self.orange = RGBColor(234, 88, 12)
        self.purple = RGBColor(124, 58, 237)
        self.red = RGBColor(220, 38, 38)

    @staticmethod
    def e(value: float):
        return Inches(value)

    def slide(self, number: int | None = None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.bg
        top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.e(self.W), self.e(0.12))
        top.fill.solid()
        top.fill.fore_color.rgb = self.cyan
        top.line.fill.background()
        if number is not None:
            self.text_box(
                slide, f'{number}/12', 14.86, 8.3, 0.65, 0.3, size=12, color=self.muted, align='right'
            )
        return slide

    def text_box(
        self,
        slide,
        text: str,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        size: float = 16,
        color: RGBColor | None = None,
        bold: bool = False,
        align: str = 'left',
        font: str = 'Arial',
    ):
        box = slide.shapes.add_textbox(self.e(x), self.e(y), self.e(w), self.e(h))
        frame = box.text_frame
        frame.clear()
        frame.margin_left = self.e(0.02)
        frame.margin_right = self.e(0.02)
        frame.margin_top = self.e(0.02)
        frame.margin_bottom = self.e(0.02)
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.TOP
        paragraphs = text.split('\n') or ['']
        for idx, chunk in enumerate(paragraphs):
            p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[align]
            run = p.add_run()
            run.text = chunk
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color or self.text
        return box

    def title(self, slide, title: str, subtitle: str, number: int) -> None:
        self.text_box(slide, title, 0.65, 0.43, 9.5, 0.45, size=27, color=self.navy, bold=True)
        if subtitle:
            self.text_box(slide, subtitle, 0.68, 0.93, 11.7, 0.35, size=14.5, color=self.muted)
        self.text_box(slide, f'{number}/12', 14.86, 0.55, 0.65, 0.3, size=12, color=self.muted, align='right')

    def panel(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        accent: RGBColor | None = None,
        fill: RGBColor | None = None,
    ):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, self.e(x), self.e(y), self.e(w), self.e(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill or self.white
        shape.line.color.rgb = self.line
        shape.line.width = Pt(1.15)
        if accent:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.e(x), self.e(y), self.e(0.08), self.e(h))
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent
            bar.line.fill.background()
        return shape

    def card(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        title: str,
        body: str,
        accent: RGBColor,
        *,
        title_size: float = 16,
        body_size: float = 13.8,
    ) -> None:
        self.panel(slide, x, y, w, h, accent)
        self.text_box(
            slide, title, x + 0.24, y + 0.16, w - 0.38, 0.34, size=title_size, color=self.navy, bold=True
        )
        self.text_box(slide, body, x + 0.24, y + 0.63, w - 0.42, h - 0.72, size=body_size, color=self.text)

    def metric(
        self, slide, x: float, y: float, w: float, h: float, value: str, label: str, accent: RGBColor
    ) -> None:
        self.panel(slide, x, y, w, h, accent)
        self.text_box(
            slide, value, x + 0.06, y + 0.12, w - 0.12, 0.37, size=22, color=accent, bold=True, align='center'
        )
        self.text_box(
            slide, label, x + 0.07, y + 0.57, w - 0.14, 0.38, size=12.2, color=self.muted, align='center'
        )

    def image(self, slide, path: Path, x: float, y: float, w: float, h: float, caption: str = '') -> None:
        self.panel(slide, x, y, w, h + (0.34 if caption else 0), None)
        with Image.open(path) as im:
            iw, ih = im.size
        scale = min((w - 0.22) / iw, (h - 0.22) / ih)
        pw = iw * scale
        ph = ih * scale
        px = x + (w - pw) / 2
        py = y + (h - ph) / 2
        slide.shapes.add_picture(str(path), self.e(px), self.e(py), width=self.e(pw), height=self.e(ph))
        if caption:
            self.text_box(
                slide,
                caption,
                x + 0.12,
                y + h + 0.06,
                w - 0.24,
                0.28,
                size=12,
                color=self.muted,
                align='center',
            )

    def bullet_panel(
        self, slide, x: float, y: float, w: float, h: float, title: str, bullets: list[str], accent: RGBColor
    ) -> None:
        self.panel(slide, x, y, w, h, accent)
        self.text_box(slide, title, x + 0.25, y + 0.17, w - 0.4, 0.35, size=16, color=self.navy, bold=True)
        bullet_text = '\n'.join(f'- {item}' for item in bullets)
        self.text_box(slide, bullet_text, x + 0.28, y + 0.68, w - 0.42, h - 0.78, size=14.2, color=self.text)

    def table(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        headers: list[str],
        rows: list[tuple[str, ...]],
        widths: list[float],
    ) -> None:
        shape = slide.shapes.add_table(
            len(rows) + 1, len(headers), self.e(x), self.e(y), self.e(w), self.e(h)
        )
        table = shape.table
        for i, width in enumerate(widths):
            table.columns[i].width = self.e(width)
        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.navy
            cell.text = header
        for row_idx, row in enumerate(rows, start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_idx % 2 else RGBColor(241, 246, 251)
                cell.text = value
        for row_idx, row in enumerate(table.rows):
            for cell in row.cells:
                cell.margin_left = self.e(0.07)
                cell.margin_right = self.e(0.07)
                cell.margin_top = self.e(0.04)
                cell.margin_bottom = self.e(0.04)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.size = Pt(13.2)
                        run.font.color.rgb = self.white if row_idx == 0 else self.text
        return shape


def build_presentation() -> Path:
    deck = CleanDeck()
    catalog = ASSETS / 'catalog_cards_crop.jpg'
    lab = ASSETS / 'engineering_lab_crop.jpg'
    cad = ASSETS / 'cad_editor_crop.jpg'
    sim = ASSETS / 'simulation_ac_crop.jpg'

    # 1
    s = deck.slide(1)
    deck.text_box(s, 'Дипломная работа', 0.72, 0.68, 6.7, 0.6, size=34, color=deck.navy, bold=True)
    deck.text_box(s, TITLE, 0.74, 1.55, 6.85, 1.55, size=19.2, color=deck.text, bold=True)
    deck.card(
        s,
        0.78,
        4.95,
        5.75,
        1.55,
        'Автор и руководитель',
        f'{AUTHOR}\nНаучный руководитель: {SUPERVISOR}\nГод защиты: 2026',
        deck.cyan,
        title_size=15,
        body_size=13.4,
    )
    deck.text_box(
        s,
        'Каталог -> схема -> расчет -> симуляция -> review -> заказ',
        0.8,
        7.25,
        6.6,
        0.35,
        size=14,
        color=deck.cyan,
        bold=True,
    )
    deck.image(s, catalog, 8.0, 0.88, 7.1, 6.65, 'Актуальный каталог и карточки компонентов')

    # 2
    s = deck.slide()
    deck.title(s, 'Актуальность и цель', 'Единый контур вместо ручного переноса данных между системами', 2)
    deck.card(
        s,
        0.72,
        1.55,
        4.35,
        2.25,
        'Проблема',
        'Каталог, datasheet, CAD, SPICE, BOM и заказ часто существуют отдельно. Пользователь вручную переносит номиналы, модели и результаты.',
        deck.blue,
        body_size=14.2,
    )
    deck.card(
        s,
        5.42,
        1.55,
        4.35,
        2.25,
        'Цель работы',
        'Создать веб-приложение, где компонент можно выбрать, проверить расчетом, применить в схеме, смоделировать и оформить в заказ.',
        deck.orange,
        body_size=14.2,
    )
    deck.card(
        s,
        10.12,
        1.55,
        4.95,
        2.25,
        'Результат',
        'DOLG связывает каталог, CAD/SIM, лабораторию, практикум, Engineering Review, BOM и покупательский контур.',
        deck.green,
        body_size=14.2,
    )
    metrics = [
        ('89', 'товаров'),
        ('43', 'РЭБ'),
        ('21', 'статья'),
        ('50', 'материалов'),
        ('13', 'уроков'),
        ('29', 'заданий'),
    ]
    for i, (value, label) in enumerate(metrics):
        deck.metric(
            s,
            1.0 + i * 2.37,
            4.65,
            1.65,
            1.12,
            value,
            label,
            [deck.cyan, deck.blue, deck.green, deck.orange, deck.purple, deck.red][i],
        )
    deck.card(
        s,
        1.25,
        6.7,
        13.45,
        0.96,
        'Формулировка эффекта',
        'Компонент, схема, расчет, измерение, учебное задание и заказ работают с одной моделью данных.',
        deck.cyan,
        title_size=15,
        body_size=14,
    )

    # 3
    s = deck.slide()
    deck.title(
        s, 'Анализ решений', 'DOLG занимает нишу между интернет-магазином, учебной средой и CAD/SIM', 3
    )
    deck.table(
        s,
        0.75,
        1.45,
        14.45,
        4.2,
        ['Класс решений', 'Сильная сторона', 'Ограничение'],
        [
            ('Маркетплейсы компонентов', 'поиск, цена, заказ', 'нет проверки схемы и расчетов'),
            ('Онлайн-симуляторы', 'быстрый эксперимент', 'слабая связь с реальными товарами'),
            ('KiCad / Altium', 'полный EDA workflow', 'высокий порог входа'),
            ('AI-CAD сервисы', 'подсказки поверх схемы', 'не всегда объяснимый источник вывода'),
            ('DOLG', 'каталог + схема + расчет + обучение', 'не заменяет промышленный PCB CAD'),
        ],
        [3.9, 4.75, 5.8],
    )
    deck.card(
        s,
        1.05,
        6.3,
        6.75,
        1.25,
        'Позиционирование',
        'Учебно-инженерный веб-сервис: пользователь проходит путь от выбора компонента до проверки и заказа без разрыва контекста.',
        deck.cyan,
        body_size=14.4,
    )
    deck.card(
        s,
        8.2,
        6.3,
        6.75,
        1.25,
        'Развитие относительно аналогов',
        'Главный следующий слой - Engineering Review, который объясняет ошибки через правила, граф схемы, BOM, расчеты и измерения.',
        deck.green,
        body_size=14.4,
    )

    # 4
    s = deck.slide()
    deck.title(s, 'Целевая аудитория', 'Сценарии использования системы', 4)
    deck.card(
        s,
        0.82,
        1.42,
        4.25,
        2.05,
        'Студенты',
        'Осваивают схемотехнику, проверяют учебные схемы, видят расчет и объяснение ошибки.',
        deck.blue,
        body_size=14.4,
    )
    deck.card(
        s,
        5.52,
        1.42,
        4.25,
        2.05,
        'Радиолюбители',
        'Подбирают компоненты, собирают простые узлы, получают набор деталей для сборки.',
        deck.green,
        body_size=14.4,
    )
    deck.card(
        s,
        10.22,
        1.42,
        4.25,
        2.05,
        'Инженеры',
        'Сравнивают номиналы, проверяют тепловой запас, повторяют удачные решения.',
        deck.orange,
        body_size=14.4,
    )
    steps = [
        ('1', 'Найти компонент'),
        ('2', 'Добавить в схему'),
        ('3', 'Запустить расчет'),
        ('4', 'Сравнить результат'),
        ('5', 'Сформировать заказ'),
    ]
    for i, (num, label) in enumerate(steps):
        x = 1.1 + i * 2.88
        deck.metric(s, x, 5.15, 1.25, 1.05, num, label, deck.cyan)
        if i < len(steps) - 1:
            line = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, deck.e(x + 1.28), deck.e(5.64), deck.e(1.38), deck.e(0.05)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = deck.line
            line.line.fill.background()
    deck.text_box(
        s, 'Типовой пользовательский маршрут', 0.9, 4.42, 6.7, 0.35, size=17, color=deck.navy, bold=True
    )

    # 5
    s = deck.slide()
    deck.title(
        s,
        'Архитектура системы',
        'Приложения Django связаны service-layer, а не дублированием формул в шаблонах',
        5,
    )
    layers = [
        (
            'UI и клиентский слой',
            'каталог, карточки, CAD-редактор, симулятор, лаборатория, обучение',
            deck.cyan,
        ),
        ('Django apps', 'shop, accounts, orders, knowledge, Dolg_APP, административный контур', deck.blue),
        (
            'Service-layer',
            'engineering_lab, project_review, learning_grader, rule_ai, cad_import, formula_steps',
            deck.green,
        ),
        (
            'Expert-first core',
            'jsonschema, rule-engine, Pint, Lark, Z3, scikit-fuzzy, graph/formula services',
            deck.orange,
        ),
        (
            'Данные',
            'SQLite demo DB, товары, РЭБ-компоненты, схемы, уроки, попытки, review snapshots',
            deck.purple,
        ),
    ]
    for i, (title, body, color) in enumerate(layers):
        deck.card(s, 0.92, 1.35 + i * 1.18, 14.1, 0.92, title, body, color, title_size=15.2, body_size=13.2)
    deck.card(
        s,
        1.35,
        7.35,
        13.2,
        0.7,
        'Ключевой принцип',
        'Формулы, проверки, импорт, review и подсказки используются повторно в лаборатории, обучении, симуляторе и AI-помощнике.',
        deck.cyan,
        title_size=14,
        body_size=13,
    )

    # 6
    s = deck.slide()
    deck.title(
        s,
        'Ключевые фрагменты реализации',
        'На слайд вынесены структуры, которые объясняют инженерный вывод',
        6,
    )
    code_font = 'Consolas'
    blocks = [
        (
            'Expert finding',
            'rule_id: led_power_margin\nseverity: risk\nevidence: {power, limit, margin}\nrecommendation: increase resistor power rating',
            deck.blue,
        ),
        (
            'Unit-safe parsing',
            "parse('10 kOhm') -> 10000 Ohm\nparse('100 nF') -> 1e-7 F\nwarning: maybe 10 instead of 10k",
            deck.green,
        ),
        (
            'Constraint solver',
            'given Vin, Vf, Iled\nfind R in E24\ncheck power and tolerance\nreturn valid variants',
            deck.orange,
        ),
    ]
    for i, (title, code, color) in enumerate(blocks):
        x = 0.9 + i * 5.0
        deck.panel(s, x, 1.55, 4.55, 4.85, color)
        deck.text_box(s, title, x + 0.25, 1.78, 4.0, 0.35, size=16, color=deck.navy, bold=True)
        deck.text_box(s, code, x + 0.28, 2.45, 4.0, 2.3, size=14.1, color=deck.text, font=code_font)
        deck.text_box(
            s,
            'Используется в review, обучении и подсказках.',
            x + 0.28,
            5.45,
            3.9,
            0.45,
            size=12.7,
            color=deck.muted,
        )
    deck.card(
        s,
        1.1,
        7.05,
        13.6,
        0.72,
        'Почему это важно',
        'Вывод системы можно объяснить: правило сработало по конкретным фактам, единицам измерения, ограничениям и данным проекта.',
        deck.cyan,
        title_size=14.2,
        body_size=13.2,
    )

    # 7
    s = deck.slide()
    deck.title(s, 'Реализованные модули', 'Функции сгруппированы по инженерному маршруту пользователя', 7)
    cards = [
        (
            'Каталог',
            '89 товаров, 43 РЭБ-компонента, datasheet, категории, no-Wikimedia media-policy',
            deck.cyan,
        ),
        ('CAD/SIM', 'редактор схем, GND/DRC, DC/AC/TRAN, demo-схемы, графики', deck.blue),
        ('Лаборатория', 'транзисторный ключ, NE555, стабилизатор, тепло, антидребезг', deck.green),
        ('Engineering Review', 'Design Health Score, BOM-риски, рекомендации, отчеты', deck.orange),
        ('Обучение', '4 маршрута, 13 уроков, 29 заданий: расчет, схема, измерение', deck.purple),
        ('AI-помощник', 'самописный rule_ai на данных проекта, без обязательной внешней LLM', deck.red),
    ]
    for i, (title, body, color) in enumerate(cards):
        x = 0.85 + (i % 3) * 5.0
        y = 1.55 + (i // 3) * 2.45
        deck.card(s, x, y, 4.35, 1.85, title, body, color, title_size=15.5, body_size=13.2)
    deck.card(
        s,
        1.2,
        7.0,
        13.2,
        0.78,
        'Итог по функциональности',
        'Сайт работает не только как магазин, а как связанный учебно-инженерный контур.',
        deck.cyan,
        title_size=14.2,
        body_size=13.5,
    )

    # 8
    s = deck.slide()
    deck.title(s, 'CAD-редактор', 'Отдельный режим проектирования схем и подготовки данных для симуляции', 8)
    deck.image(s, cad, 0.75, 1.42, 9.95, 5.95, 'Режим CAD: схема, элементы, соединения и рабочая область')
    deck.bullet_panel(
        s,
        11.05,
        1.45,
        4.1,
        1.95,
        'Что проверяется',
        ['наличие GND', 'источник питания', 'соединения', 'номиналы'],
        deck.blue,
    )
    deck.bullet_panel(
        s,
        11.05,
        3.8,
        4.1,
        1.65,
        'Что формируется',
        ['scheme_data', 'BOM', 'netlist', 'review input'],
        deck.green,
    )
    deck.bullet_panel(
        s, 11.05, 5.88, 4.1, 1.48, 'Зачем это нужно', ['не просто рисунок', 'данные для анализа'], deck.orange
    )

    # 9
    s = deck.slide()
    deck.title(
        s,
        'Симуляция и измерения',
        'Расчет, графики и метрики используются в лаборатории, review и обучении',
        9,
    )
    deck.image(s, sim, 0.75, 1.42, 9.15, 5.95, 'Графики и результаты AC/DC/TRAN анализа')
    deck.card(
        s,
        10.25,
        1.45,
        4.8,
        1.58,
        'Режимы анализа',
        'DC / OP, AC, TRAN, Bode plot, FFT spectrum, Monte Carlo tolerance.',
        deck.blue,
        body_size=13.8,
    )
    deck.card(
        s,
        10.25,
        3.33,
        4.8,
        1.58,
        'Метрики',
        'Напряжение узла, ток ветви, RMS, частота, duty cycle, мощность, температура.',
        deck.green,
        body_size=13.8,
    )
    deck.card(
        s,
        10.25,
        5.21,
        4.8,
        1.58,
        'Инженерная оценка',
        'Норма, риск, перегрев, нужен запас. Результат не остается просто числом.',
        deck.orange,
        body_size=13.8,
    )
    deck.text_box(
        s,
        'Expected vs measured связывает лабораторный расчет и фактическое измерение из симуляции.',
        1.0,
        7.55,
        13.9,
        0.35,
        size=14,
        color=deck.cyan,
        bold=True,
        align='center',
    )

    # 10
    s = deck.slide()
    deck.title(
        s, 'Проверка результата', 'Тесты и демонстрационные проверки подтверждают готовность к защите', 10
    )
    deck.image(s, lab, 0.75, 1.45, 7.15, 5.5, 'Лаборатория: расчет и оценка результата')
    deck.table(
        s,
        8.25,
        1.45,
        6.65,
        3.75,
        ['Проверка', 'Результат'],
        [
            ('manage.py check', '0 issues'),
            ('makemigrations', 'No changes detected'),
            ('check_demo_ready', 'OK, expert_stack'),
            ('check_data_integrity', 'OK, 0 errors'),
            ('Targeted tests', '18 expert / 16 learning / 8 search'),
        ],
        [2.6, 4.05],
    )
    for i, (value, label, color) in enumerate(
        [
            ('OK', 'demo-ready', deck.green),
            ('OK', 'integrity', deck.green),
            ('0', 'warnings', deck.blue),
            ('89', 'products', deck.cyan),
        ]
    ):
        deck.metric(s, 8.45 + i * 1.55, 5.92, 1.22, 1.0, value, label, color)
    deck.card(
        s,
        8.25,
        7.2,
        6.65,
        0.68,
        'Вывод',
        'Проверяется не только Django-конфигурация, но и наполнение, demo-сценарии, обучение, expert stack и media-policy.',
        deck.cyan,
        title_size=13.8,
        body_size=12.8,
    )

    # 11
    s = deck.slide()
    deck.title(s, 'План развития', 'Новая область вводится пакетом: фича + обучалка + документация', 11)
    roadmap = [
        ('1', 'Expert Review Core', 'расширение rule packs, explainable findings, Learning-by-review'),
        ('2', 'Measurement Core', 'probes, expected vs measured, sweep, сохранение измерений'),
        ('3', 'CAD Import', 'LTspice/KiCad subset, import preview, review после импорта'),
        ('4', 'Комментарии', 'обсуждения к товарам, урокам, статьям, demo-схемам и review'),
        ('5', 'Neural deep analysis', 'PyTorch/GOLEM позже, только поверх expert baseline'),
    ]
    for i, (num, title, body) in enumerate(roadmap):
        y = 1.38 + i * 1.18
        deck.metric(s, 0.95, y, 0.88, 0.75, num, 'этап', deck.cyan)
        deck.card(
            s,
            2.05,
            y,
            12.75,
            0.75,
            title,
            body,
            [deck.blue, deck.green, deck.orange, deck.purple, deck.red][i],
            title_size=14.4,
            body_size=12.6,
        )
    deck.card(
        s,
        1.15,
        7.35,
        13.4,
        0.68,
        'Принцип развития',
        'Сначала объяснимые экспертные системы и ограничения, затем оптимизация, и только после этого нейронные подсказки.',
        deck.cyan,
        title_size=13.8,
        body_size=13,
    )

    # 12
    s = deck.slide(12)
    deck.text_box(s, 'Спасибо за внимание!', 0.9, 1.02, 8.2, 0.75, size=36, color=deck.navy, bold=True)
    deck.text_box(
        s,
        'DOLG объединяет каталог, CAD/SIM, инженерную лабораторию, обучение и экспертную проверку проекта.',
        0.94,
        2.15,
        7.45,
        0.95,
        size=18,
        color=deck.text,
    )
    deck.card(
        s,
        1.0,
        4.15,
        5.9,
        1.65,
        'Контактная информация',
        f'{AUTHOR}\nEmail: buryako@internet.com\nТелефон: +7 (903) 439-44-87',
        deck.cyan,
        title_size=15,
        body_size=13.4,
    )
    deck.card(
        s,
        7.55,
        4.15,
        6.9,
        1.65,
        'Материалы к защите',
        'Диплом, речь и презентация актуализированы. Проверки проекта проходят, текущие данные зафиксированы в документах.',
        deck.green,
        title_size=15,
        body_size=13.4,
    )
    deck.text_box(s, 'Вопросы?', 5.35, 7.1, 5.2, 0.55, size=28, color=deck.cyan, bold=True, align='center')

    docs_pptx = first_file('Презентация_DOLG_финальная_*.pptx')
    actual_pptx = DOCS / 'Презентация_DOLG_актуальная_20260519.pptx'
    white_pptx = DOCS / 'Презентация_DOLG_белая_читабельная_20260519.pptx'
    deck.prs.save(str(docs_pptx))
    deck.prs.save(str(actual_pptx))
    deck.prs.save(str(white_pptx))

    download_target = (
        DOWNLOADS
        / 'Razrabotka-veb-prilozheniya-dlya-prodazhi-radio-i-elektronnyh-komponentov-so-vstroennymi-instrumenta_updated_20260518.pptx'
    )
    download_white = (
        DOWNLOADS
        / 'Razrabotka-veb-prilozheniya-dlya-prodazhi-radio-i-elektronnyh-komponentov-so-vstroennymi-instrumenta_fixed_white_20260519.pptx'
    )
    if download_target.exists():
        copy_file_safe(docs_pptx, download_target)
    copy_file_safe(docs_pptx, download_white)

    Presentation(str(docs_pptx))
    Presentation(str(actual_pptx))
    Presentation(str(white_pptx))
    return white_pptx


def main() -> None:
    backup_dir = DOCS / 'presentation_backups' / '20260519_white_before_script'
    backup_dir.mkdir(parents=True, exist_ok=True)
    for pattern in [
        'Презентация_DOLG_*.pptx',
        'Речь_и_вопросы_к_защите_DOLG_*.md',
        'Речь_и_вопросы_к_защите_DOLG_*.docx',
    ]:
        for path in DOCS.glob(pattern):
            if not path.name.startswith('~$'):
                backup_file(path, backup_dir)
    download_pptx = (
        DOWNLOADS
        / 'Razrabotka-veb-prilozheniya-dlya-prodazhi-radio-i-elektronnyh-компонентов-so-vstroennymi-instrumenta_updated_20260518.pptx'
    )
    backup_file(download_pptx, backup_dir)
    download_pptx = (
        DOWNLOADS
        / 'Razrabotka-veb-prilozheniya-dlya-prodazhi-radio-i-elektronnyh-komponentov-so-vstroennymi-instrumenta_updated_20260518.pptx'
    )
    backup_file(download_pptx, backup_dir)

    update_speech_intro()
    out = build_presentation()
    print(f'Built clean white presentation: {out}')


if __name__ == '__main__':
    main()
