from pathlib import Path
from shutil import copy2

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / 'docs' / 'Презентация_DOLG_финальная_20260513_v5.pptx'
RELEASE_OUT = ROOT / 'release' / 'DOLG_final_20260513_v5' / 'Презентация_DOLG_финальная_20260513_v5.pptx'
IMG_DIR = ROOT / 'docs' / 'diploma_assets' / 'screenshots'
V6 = IMG_DIR / 'presentation_v6'
CROP = V6 / 'crops'
CROP.mkdir(parents=True, exist_ok=True)

TOPIC = (
    'Разработка веб-приложения для продажи радио- и электронных компонентов '
    'со встроенными инструментами проектирования и симуляции схем.'
)
AUTHOR = 'Буряко Дмитрий Сергеевич'
SUPERVISOR = 'Буланов Сергей Георгиевич'
YEAR = '2026'
TOTAL = 12

BG = RGBColor(245, 248, 251)
DARK = RGBColor(18, 34, 61)
MUTED = RGBColor(78, 96, 124)
CYAN = RGBColor(0, 157, 181)
BLUE = RGBColor(47, 111, 225)
GREEN = RGBColor(31, 178, 112)
ORANGE = RGBColor(236, 87, 25)
PURPLE = RGBColor(115, 75, 230)
RED = RGBColor(220, 48, 55)
LIGHT_CYAN = RGBColor(226, 249, 252)
CARD = RGBColor(255, 255, 255)
BORDER = RGBColor(204, 215, 229)
CODE_BG = RGBColor(15, 29, 50)
CODE_TEXT = RGBColor(229, 241, 255)

W = Inches(13.333333)
H = Inches(7.5)


def crop_image(src, dst, box):
    image = Image.open(src).convert('RGB')
    image = image.crop(box)
    image.save(dst, quality=92)
    return dst


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color=None, width=1, transparency=0):
    if color is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def text_box(
    slide,
    x,
    y,
    w,
    h,
    text='',
    size=16,
    color=DARK,
    bold=False,
    align=PP_ALIGN.LEFT,
    font='Arial',
    valign=MSO_ANCHOR.TOP,
    margin=0.04,
    italic=False,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_bg(slide, top=True, left=False):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    set_fill(bg, BG)
    set_line(bg, None)
    if top:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.17))
        set_fill(bar, CYAN)
        set_line(bar, None)
    if left:
        lbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), H)
        set_fill(lbar, CYAN)
        set_line(lbar, None)


def footer(slide, idx):
    text_box(
        slide,
        Inches(0.65),
        Inches(7.05),
        Inches(3.7),
        Inches(0.24),
        'DOLG · дипломная работа',
        7.5,
        MUTED,
    )
    text_box(
        slide,
        Inches(11.55),
        Inches(7.05),
        Inches(0.85),
        Inches(0.24),
        f'{idx}/{TOTAL}',
        8,
        MUTED,
        align=PP_ALIGN.RIGHT,
    )


def title(slide, text, subtitle=None, idx=None):
    text_box(slide, Inches(0.65), Inches(0.38), Inches(8.6), Inches(0.5), text, 24, DARK, True)
    if subtitle:
        text_box(
            slide,
            Inches(0.65),
            Inches(0.88),
            Inches(10.4),
            Inches(0.35),
            subtitle,
            12.5,
            MUTED,
        )
    if idx:
        footer(slide, idx)


def card(slide, x, y, w, h, heading, body, accent=CYAN, heading_size=14, body_size=10.5):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(rect, CARD)
    set_line(rect, BORDER, 1)
    rect.shadow.inherit = False
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), h)
    set_fill(stripe, accent)
    set_line(stripe, None)
    text_box(
        slide,
        x + Inches(0.23),
        y + Inches(0.16),
        w - Inches(0.42),
        Inches(0.35),
        heading,
        heading_size,
        DARK,
        True,
    )
    text_box(
        slide,
        x + Inches(0.23),
        y + Inches(0.6),
        w - Inches(0.42),
        h - Inches(0.68),
        body,
        body_size,
        MUTED,
    )
    return rect


def chip(slide, x, y, w, label, value, color=CYAN):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.56))
    set_fill(rect, CARD)
    set_line(rect, color, 1.1)
    text_box(
        slide,
        x + Inches(0.08),
        y + Inches(0.08),
        w - Inches(0.16),
        Inches(0.18),
        label,
        7.5,
        MUTED,
        align=PP_ALIGN.CENTER,
    )
    text_box(
        slide,
        x + Inches(0.08),
        y + Inches(0.25),
        w - Inches(0.16),
        Inches(0.24),
        value,
        12,
        color,
        True,
        align=PP_ALIGN.CENTER,
    )


def image_panel(slide, path, x, y, w, h):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(panel, RGBColor(255, 255, 255))
    set_line(panel, BORDER, 1)
    slide.shapes.add_picture(
        str(path),
        x + Inches(0.06),
        y + Inches(0.06),
        width=w - Inches(0.12),
        height=h - Inches(0.12),
    )


def code_card(slide, x, y, w, h, heading, code, accent=CYAN):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(rect, CODE_BG)
    set_line(rect, accent, 1.1)
    text_box(
        slide,
        x + Inches(0.15),
        y + Inches(0.12),
        w - Inches(0.3),
        Inches(0.28),
        heading,
        10.5,
        RGBColor(181, 232, 255),
        True,
    )
    box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.48), w - Inches(0.3), h - Inches(0.58))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    run = frame.paragraphs[0].add_run()
    run.text = code
    run.font.name = 'Consolas'
    run.font.size = Pt(7.5)
    run.font.color.rgb = CODE_TEXT


def make_table(slide, x, y, w, h, data, font_size=8.5):
    rows, cols = len(data), len(data[0])
    shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = shape.table
    for row_idx, row in enumerate(data):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = LIGHT_CYAN if row_idx == 0 else CARD
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = 'Arial'
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = DARK if row_idx == 0 else MUTED
                    run.font.bold = row_idx == 0


def connector(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = CYAN
    line.line.width = Pt(2)


def build():
    catalog_img = crop_image(
        V6 / '01_catalog_cards_current.png',
        CROP / 'catalog_cards_crop.jpg',
        (120, 285, 1360, 890),
    )
    lab_img = crop_image(
        V6 / '02_engineering_lab_results.png',
        CROP / 'engineering_lab_crop.jpg',
        (120, 175, 1325, 875),
    )
    cad_img = crop_image(
        V6 / '03_cad_editor_current.png',
        CROP / 'cad_editor_crop.jpg',
        (0, 150, 1440, 890),
    )
    sim_img = crop_image(
        IMG_DIR / '02_simulation_ac_graph.png',
        CROP / 'simulation_ac_crop.jpg',
        (10, 165, 1665, 1065),
    )

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # 1. Title
    slide = prs.slides.add_slide(blank)
    add_bg(slide, top=False, left=True)
    text_box(slide, Inches(0.72), Inches(0.72), Inches(5.1), Inches(0.55), 'Дипломная работа', 30, DARK, True)
    text_box(slide, Inches(0.72), Inches(1.5), Inches(5.45), Inches(1.35), TOPIC, 20, DARK, True)
    text_box(
        slide,
        Inches(0.72),
        Inches(5.45),
        Inches(5.35),
        Inches(0.92),
        f'Автор: {AUTHOR}\nНаучный руководитель: {SUPERVISOR}\nГод защиты: {YEAR}',
        11,
        MUTED,
    )
    image_panel(slide, catalog_img, Inches(6.15), Inches(0.75), Inches(6.25), Inches(5.7))

    # 2. Relevance and goal
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(slide, 'Актуальность и цель', 'Рынок, ограничения существующих инструментов и цель ВКР', 2)
    card(
        slide,
        Inches(0.8),
        Inches(1.55),
        Inches(3.55),
        Inches(3.45),
        'Рынок и EDA-среда',
        'Рост спроса на электронные компоненты и инструменты проектирования усиливает потребность в веб-сервисах, где подбор, расчёт и документация находятся в одном рабочем маршруте.',
        BLUE,
        13.5,
        11,
    )
    card(
        slide,
        Inches(4.9),
        Inches(1.55),
        Inches(3.55),
        Inches(3.45),
        'Практический разрыв',
        'Каталог, CAD и SPICE-среда часто используются отдельно. Номиналы, спецификации, BOM и результаты расчётов переносятся между системами вручную.',
        ORANGE,
        13.5,
        11,
    )
    card(
        slide,
        Inches(9.0),
        Inches(1.55),
        Inches(3.55),
        Inches(3.45),
        'Цель работы',
        'Разработать веб-приложение для подбора, приобретения радио- и электронных компонентов, полного проектирования и симуляции электронной схемы в онлайн-режиме.',
        GREEN,
        13.5,
        11,
    )
    chip(slide, Inches(1.15), Inches(5.55), Inches(2.2), 'реализовано', 'каталог + заказ', BLUE)
    chip(slide, Inches(3.95), Inches(5.55), Inches(2.2), 'реализовано', 'CAD/SIM', CYAN)
    chip(slide, Inches(6.75), Inches(5.55), Inches(2.2), 'добавлено', 'лаборатория', GREEN)
    chip(slide, Inches(9.55), Inches(5.55), Inches(2.2), 'добавлено', 'обучение', PURPLE)

    # 3. Solutions analysis
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(
        slide,
        'Анализ решений',
        'Позиционирование проекта относительно магазинов, онлайн-симуляторов и CAD-систем',
        3,
    )
    make_table(
        slide,
        Inches(0.65),
        Inches(1.35),
        Inches(12.05),
        Inches(3.85),
        [
            ['Класс решений', 'Сильная сторона', 'Ограничение', 'Как закрывает DOLG'],
            [
                'Интернет-магазины',
                'каталог, цена, наличие, заказ',
                'нет проверки схемы и инженерного расчёта',
                'товар связан с документацией, BOM и заказом',
            ],
            [
                'Онлайн-симуляторы',
                'быстрые учебные эксперименты',
                'абстрактные элементы без рыночной позиции',
                'схема работает рядом с каталогом компонентов',
            ],
            [
                'CAD/EDA',
                'профессиональная разработка схем и плат',
                'высокий порог входа и отдельный контур закупки',
                'упрощённый веб-маршрут для учебных и практических задач',
            ],
            [
                'DOLG',
                'каталог, знания, CAD, симуляция, лаборатория',
                'демонстрационный масштаб проекта ВКР',
                'единый маршрут: компонент -> схема -> расчёт -> заказ',
            ],
        ],
        8.4,
    )
    card(
        slide,
        Inches(1.2),
        Inches(5.42),
        Inches(10.95),
        Inches(1.0),
        'Практическая ниша',
        'Учебно-практический веб-сервис для работы с электронным компонентом: от поиска и анализа параметров до проверки поведения схемы и формирования заказа.',
        CYAN,
        11.5,
        9.2,
    )

    # 4. Audience
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(
        slide,
        'Целевая аудитория и сценарии',
        'Основные группы пользователей и рабочий маршрут внутри системы',
        4,
    )
    card(
        slide,
        Inches(0.8),
        Inches(1.25),
        Inches(4.7),
        Inches(1.12),
        'Студенты',
        'изучают схемотехнику через расчёты, готовые схемы и лабораторные задания.',
        BLUE,
        12.5,
        8.8,
    )
    card(
        slide,
        Inches(0.8),
        Inches(2.55),
        Inches(4.7),
        Inches(1.12),
        'Радиолюбители',
        'подбирают компоненты, проверяют идею и собирают список деталей для покупки.',
        GREEN,
        12.5,
        8.8,
    )
    card(
        slide,
        Inches(0.8),
        Inches(3.85),
        Inches(4.7),
        Inches(1.12),
        'Инженеры',
        'сравнивают номиналы, оценивают тепловой запас и сохраняют проверенные решения.',
        ORANGE,
        12.5,
        8.8,
    )
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.05), Inches(1.25), Inches(5.95), Inches(3.7)
    )
    set_fill(panel, CARD)
    set_line(panel, BORDER, 1)
    text_box(
        slide,
        Inches(6.35),
        Inches(1.52),
        Inches(5.25),
        Inches(0.35),
        'Типовой пользовательский маршрут',
        14,
        DARK,
        True,
    )
    steps = [
        'Найти компонент',
        'Открыть карточку и документацию',
        'Проверить расчёт в лаборатории',
        'Собрать схему и запустить симуляцию',
        'Сформировать BOM и заказ',
    ]
    for i, step in enumerate(steps):
        yy = Inches(2.05 + i * 0.48)
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.45), yy, Inches(0.18), Inches(0.18))
        set_fill(oval, CYAN)
        set_line(oval, None)
        text_box(
            slide,
            Inches(6.78),
            yy - Inches(0.04),
            Inches(4.7),
            Inches(0.26),
            f'{i + 1}. {step}',
            10.5,
            DARK if i in (0, 2) else MUTED,
            bold=i in (0, 2),
        )
        if i < len(steps) - 1:
            connector(slide, Inches(6.54), yy + Inches(0.18), Inches(6.54), yy + Inches(0.43))
    for i, (value, label) in enumerate(
        [
            ('89', 'товаров'),
            ('43', 'РЭБ'),
            ('21', 'статья'),
            ('50', 'материалов'),
            ('10', 'уроков'),
            ('23', 'задания'),
        ]
    ):
        chip(
            slide,
            Inches(0.95 + i * 1.9),
            Inches(5.45),
            Inches(1.45),
            label,
            value,
            [BLUE, CYAN, GREEN, ORANGE, PURPLE, RED][i],
        )

    # 5. Architecture
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(
        slide, 'Архитектура', 'Фактическая структура Django-проекта, сервисного слоя и клиентских модулей', 5
    )
    card(
        slide,
        Inches(0.55),
        Inches(1.25),
        Inches(2.75),
        Inches(4.65),
        'Django apps',
        'shop — каталог, поиск, сравнение, BOM\naccounts — регистрация, профиль\norders — корзина и заказы\nknowledge — энциклопедия, обучение, лаборатория\nDolg_APP — CAD, симуляция, проекты',
        BLUE,
        12.5,
        9.2,
    )
    card(
        slide,
        Inches(3.65),
        Inches(1.25),
        Inches(2.75),
        Inches(4.65),
        'Service-layer',
        'engineering_lab.py — расчёты и оценка\nlearning_grader.py — автопроверка задач\nschematic_validation.py — DRC/проверка схем\nbom_match — подбор позиций каталога\ncheck_demo_ready — контроль демо',
        GREEN,
        12.5,
        9.2,
    )
    card(
        slide,
        Inches(6.75),
        Inches(1.25),
        Inches(2.75),
        Inches(4.65),
        'Client/UI',
        'единые карточки товара\nCanvas2D CAD с УГО и ГОСТ-шаблонами\nngspice.wasm + JS-MNA fallback\nинженерная лаборатория\nпрактикум обучения',
        ORANGE,
        12.5,
        9.2,
    )
    card(
        slide,
        Inches(9.85),
        Inches(1.25),
        Inches(2.75),
        Inches(4.65),
        'Data & quality',
        'SQLite для demo/dev\nJSONField для параметров компонентов\nmedia: фото, gif, видео, файлы\n146 тестов, 6 skipped\ndemo-ready URL smoke',
        PURPLE,
        12.5,
        9.2,
    )
    for x in [Inches(3.32), Inches(6.42), Inches(9.52)]:
        connector(slide, x, Inches(3.55), x + Inches(0.28), Inches(3.55))
    text_box(
        slide,
        Inches(1.25),
        Inches(6.15),
        Inches(10.9),
        Inches(0.42),
        'Архитектура разделяет торговую, учебную и инженерную части, а общие расчёты вынесены из шаблонов в переиспользуемые сервисы.',
        11.2,
        MUTED,
        align=PP_ALIGN.CENTER,
    )

    # 6. Code
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(
        slide,
        'Ключевые фрагменты кода',
        'Фрагменты, которые связывают лабораторию, обучение, проверку схем и demo-ready контроль',
        6,
    )
    code_card(
        slide,
        Inches(0.55),
        Inches(1.25),
        Inches(3.0),
        Inches(2.25),
        'Расчётная лаборатория',
        "def calculate_lab(kind, payload):\n    calculators = {\n        'transistor_switch': _transistor_switch,\n        'ne555_astable': _ne555_astable,\n        'linear_regulator': _linear_regulator,\n    }\n    result = calculators[kind](payload or {})\n    return _with_feedback(result)",
        BLUE,
    )
    code_card(
        slide,
        Inches(3.78),
        Inches(1.25),
        Inches(3.0),
        Inches(2.25),
        'Обучение использует те же метрики',
        "def grade_simulation_task(task, scheme_data, result):\n    rubric = task.rubric or {}\n    metric = rubric.get('metric')\n    value = extract_measurement(result, metric, rubric)\n    return evaluate_measurement(metric, value, rubric)",
        GREEN,
    )
    code_card(
        slide,
        Inches(7.0),
        Inches(1.25),
        Inches(3.0),
        Inches(2.25),
        'DRC схемы как сервис',
        "def validate_scheme_data(scheme_data):\n    components = _components(scheme_data)\n    errors, warnings = [], []\n    if not _has_ground(components):\n        errors.append('нет GND')\n    return {'errors': errors, 'warnings': warnings,\n            'metrics': metrics}",
        ORANGE,
    )
    code_card(
        slide,
        Inches(10.22),
        Inches(1.25),
        Inches(2.55),
        Inches(2.25),
        'Demo-ready контроль',
        "minimums = {\n  'products': 80,\n  'learning_tracks': 2,\n  'learning_lessons': 10,\n  'learning_tasks': 22,\n}\nsmoke_urls += ['/knowledge/lab/']",
        PURPLE,
    )
    code_card(
        slide,
        Inches(0.95),
        Inches(4.05),
        Inches(5.45),
        Inches(1.85),
        'Переиспользование формул',
        "def lab_expected_value(config):\n    result = calculate_lab(config['kind'], config['inputs'])\n    metric = config['metric']\n    return result['outputs'][metric]['value']\n\n# одно место для лаборатории и практикума",
        CYAN,
    )
    code_card(
        slide,
        Inches(6.95),
        Inches(4.05),
        Inches(5.0),
        Inches(1.85),
        'Проверка схемного задания',
        "required_types = rubric.get('required_component_types', [])\nnominals = rubric.get('required_nominal_ranges', [])\nconnections = _connection_graph(scheme_data)\n# feedback возвращается в LearningAttempt",
        RED,
    )

    # 7. Implemented modules
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(
        slide,
        'Реализованные модули',
        'Функции продукта и свойства данных, которые поддерживают пользовательский маршрут',
        7,
    )
    card(
        slide,
        Inches(0.65),
        Inches(1.22),
        Inches(3.65),
        Inches(1.38),
        'Каталог и подбор',
        'поиск, фильтры по параметрам, наличие, производитель, lifecycle; единое представление компонента в разных местах интерфейса.',
        BLUE,
        12.5,
        10.2,
    )
    card(
        slide,
        Inches(4.85),
        Inches(1.22),
        Inches(3.65),
        Inches(1.38),
        'Корзина и заказ',
        'BOM из схемы может перейти в корзину; поддержаны оформление, повтор заказа и проверка остатков.',
        GREEN,
        12.5,
        10.2,
    )
    card(
        slide,
        Inches(9.05),
        Inches(1.22),
        Inches(3.65),
        Inches(1.38),
        'Энциклопедия',
        '21 статья и 50 материалов: внутренние ссылки, изображения, видео, файлы и связанные товары.',
        ORANGE,
        12.5,
        10.2,
    )
    card(
        slide,
        Inches(0.65),
        Inches(3.22),
        Inches(3.65),
        Inches(1.38),
        'Инженерная лаборатория',
        'расчёты транзисторного ключа, NE555, стабилизатора, RC-антидребезга и теплового запаса.',
        CYAN,
        12.5,
        10.2,
    )
    card(
        slide,
        Inches(4.85),
        Inches(3.22),
        Inches(3.65),
        Inches(1.38),
        'Практикум обучения',
        '2 маршрута, 10 уроков, 23 задания: численные ответы, сборка схемы и измерение результата.',
        PURPLE,
        12.5,
        10.2,
    )
    card(
        slide,
        Inches(9.05),
        Inches(3.22),
        Inches(3.65),
        Inches(1.38),
        'Проекты и спецификация',
        'сохранение схем, история запусков, PDF/XLSX-выгрузки, demo-проекты и переход к спецификации.',
        RED,
        12.5,
        10.2,
    )
    text_box(
        slide,
        Inches(0.95),
        Inches(5.35),
        Inches(11.35),
        Inches(0.48),
        'Ассортимент и карточки здесь рассматриваются как свойства каталога, а не как отдельные функции: они обеспечивают работу модулей подбора, расчёта и заказа.',
        11,
        MUTED,
        align=PP_ALIGN.CENTER,
    )

    # 8. CAD
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(slide, 'CAD-редактор', 'Отдельный режим для чертежей, ГОСТ-шаблонов, компонентов и экспорта', 8)
    image_panel(slide, cad_img, Inches(0.55), Inches(1.2), Inches(8.65), Inches(5.45))
    card(
        slide,
        Inches(9.55),
        Inches(1.25),
        Inches(2.75),
        Inches(1.05),
        'Режим CAD',
        'сетка, snap, слои, свойства, инструментальные кнопки и отдельный холст.',
        CYAN,
        12,
        9.2,
    )
    card(
        slide,
        Inches(9.55),
        Inches(2.65),
        Inches(2.75),
        Inches(1.05),
        'ГОСТ и чертежи',
        'рамки, штампы, тест-сценарии, размерные линии и выноски.',
        BLUE,
        12,
        9.2,
    )
    card(
        slide,
        Inches(9.55),
        Inches(4.05),
        Inches(2.75),
        Inches(1.05),
        'Экспорт',
        'PDF A3, PNG, SVG, JSON и передача в симулятор.',
        GREEN,
        12,
        9.2,
    )

    # 9. Simulation
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(slide, 'Симуляция и измерения', 'Схема, SPICE-анализ, BOM и инженерная обратная связь', 9)
    image_panel(slide, sim_img, Inches(0.55), Inches(1.2), Inches(8.1), Inches(5.45))
    card(
        slide,
        Inches(9.05),
        Inches(1.25),
        Inches(3.25),
        Inches(1.0),
        'Режимы анализа',
        'DC, AC и TRAN; при недоступности WASM используется JS-MNA fallback.',
        BLUE,
        12,
        9.3,
    )
    card(
        slide,
        Inches(9.05),
        Inches(2.55),
        Inches(3.25),
        Inches(1.0),
        'Метрики',
        'напряжение узла, ток ветви, RMS, частота, duty cycle, мощность элемента.',
        GREEN,
        12,
        9.3,
    )
    card(
        slide,
        Inches(9.05),
        Inches(3.85),
        Inches(3.25),
        Inches(1.0),
        'Связь с каталогом',
        'из схемы формируется BOM, затем позиции сопоставляются с товарами.',
        ORANGE,
        12,
        9.3,
    )
    card(
        slide,
        Inches(9.05),
        Inches(5.15),
        Inches(3.25),
        Inches(1.0),
        'Лаборатория',
        'те же расчёты используются в обучении и проверке практических задач.',
        PURPLE,
        12,
        9.3,
    )

    # 10. Validation
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(
        slide,
        'Проверка результата',
        'Тесты, demo-ready контроль и лабораторный результат как наглядное подтверждение',
        10,
    )
    make_table(
        slide,
        Inches(0.65),
        Inches(1.35),
        Inches(4.2),
        Inches(3.0),
        [
            ['Уровень проверки', 'Что подтверждается'],
            ['Unit/model tests', 'модели, расчёты, grader, BOM, заказы'],
            ['View/API tests', 'страницы лаборатории, обучения, поиска и API'],
            ['Demo-ready', 'минимумы данных и URL smoke, включая /knowledge/lab/'],
            ['Data integrity', 'изображения, статьи, ссылки, demo-схемы'],
        ],
        8.4,
    )
    chip(slide, Inches(0.75), Inches(4.78), Inches(1.32), 'tests', '146', BLUE)
    chip(slide, Inches(2.2), Inches(4.78), Inches(1.32), 'OK', '140', GREEN)
    chip(slide, Inches(3.65), Inches(4.78), Inches(1.32), 'skipped', '6', ORANGE)
    image_panel(slide, lab_img, Inches(5.25), Inches(1.25), Inches(7.2), Inches(5.45))
    text_box(
        slide,
        Inches(0.75),
        Inches(5.78),
        Inches(4.45),
        Inches(0.45),
        'На слайде справа — результат, который пользователь получает в инженерной лаборатории: численные метрики сразу сопровождаются оценкой “норма” или “риск”.',
        9.5,
        MUTED,
    )

    # 11. Plan
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(slide, 'План развития', 'Следующий фронт работ после добавления лаборатории и практикума', 11)
    for i, (head, body) in enumerate(
        [
            (
                '1. CAD/SIM ядро',
                'probes и курсоры графиков, DRC/ERC, устойчивый netlist, улучшение smart wiring и читаемости схем',
            ),
            (
                '2. Лаборатория измерений',
                'графики измерений, branch current, RMS, duty cycle, мощность элемента, температура, журнал лабораторных запусков',
            ),
            (
                '3. Обучение пакетами',
                'каждая крупная фича добавляется вместе с учебным блоком, задачами и demo-сценарием',
            ),
            (
                '4. Интеграции CAD/EDA',
                'KiCad, LTspice, EDIF, Gerber/Excellon, расширенная библиотека SPICE-моделей',
            ),
            (
                '5. Production-контур',
                'PostgreSQL, HTTPS, static/media storage, мониторинг, резервное копирование и регламент обновлений',
            ),
        ]
    ):
        card(
            slide,
            Inches(0.9),
            Inches(1.15 + i * 1.03),
            Inches(11.45),
            Inches(0.94),
            head,
            body,
            [BLUE, CYAN, GREEN, ORANGE, PURPLE][i],
            12.2,
            9.0,
        )
    text_box(
        slide,
        Inches(1.2),
        Inches(6.55),
        Inches(10.8),
        Inches(0.32),
        'Базовое правило дальнейшего развития: крупная фича -> обучающий блок -> документация и демонстрационный сценарий.',
        10.5,
        DARK,
        True,
        align=PP_ALIGN.CENTER,
    )

    # 12. Final
    slide = prs.slides.add_slide(blank)
    add_bg(slide, top=False, left=True)
    text_box(
        slide, Inches(0.9), Inches(1.05), Inches(6.0), Inches(0.75), 'Спасибо за внимание', 34, DARK, True
    )
    text_box(slide, Inches(0.92), Inches(2.05), Inches(6.1), Inches(1.1), TOPIC, 18, DARK, True)
    card(
        slide,
        Inches(0.95),
        Inches(3.42),
        Inches(5.4),
        Inches(1.38),
        'Контактная информация',
        f'Автор: {AUTHOR}\nНаучный руководитель: {SUPERVISOR}\nГод защиты: {YEAR}',
        CYAN,
        13,
        9.6,
    )
    card(
        slide,
        Inches(0.95),
        Inches(5.08),
        Inches(5.4),
        Inches(1.18),
        'Демонстрация',
        'Локальный маршрут: /demo/\nКлючевой сценарий: каталог -> лаборатория -> CAD -> симуляция -> BOM -> заказ',
        GREEN,
        13,
        9.0,
    )
    image_panel(slide, catalog_img, Inches(7.05), Inches(1.05), Inches(5.15), Inches(4.85))
    text_box(
        slide,
        Inches(10.95),
        Inches(7.05),
        Inches(0.85),
        Inches(0.24),
        '12/12',
        8,
        MUTED,
        align=PP_ALIGN.RIGHT,
    )

    prs.save(OUT)
    if RELEASE_OUT.parent.exists():
        copy2(OUT, RELEASE_OUT)
    print(f'saved: {OUT}')
    if RELEASE_OUT.parent.exists():
        print(f'copied: {RELEASE_OUT}')


if __name__ == '__main__':
    build()
