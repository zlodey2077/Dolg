from __future__ import annotations

import shutil
from pathlib import Path

from docx import Document
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
DOWNLOADS = Path.home() / "Downloads"
ASSETS = DOCS / "diploma_assets" / "screenshots" / "presentation_v6" / "crops"

TITLE = (
    "Разработка веб-приложения для продажи радио- и электронных компонентов "
    "со встроенными инструментами проектирования и симуляции схем"
)
AUTHOR = "Буряко Дмитрий Сергеевич"
SUPERVISOR = "Буланов Сергей Георгиевич"


def first_doc(pattern: str) -> Path:
    files = [path for path in DOCS.glob(pattern) if not path.name.startswith("~$")]
    if not files:
        raise FileNotFoundError(pattern)
    return sorted(files, key=lambda path: path.stat().st_mtime, reverse=True)[0]


def backup(path: Path, backup_dir: Path) -> None:
    if path.exists():
        target = backup_dir / path.name
        if not target.exists():
            shutil.copy2(path, target)


def copy_safe(source: Path, target: Path) -> Path:
    try:
        shutil.copy2(source, target)
        return target
    except PermissionError:
        fallback = target.with_name(f"{target.stem}_final_clean_20260519{target.suffix}")
        shutil.copy2(source, fallback)
        return fallback


def update_speech_note() -> None:
    md_path = first_doc("Речь_и_вопросы_к_защите_DOLG_*.md")
    text = md_path.read_text(encoding="utf-8")
    text = text.replace(
        "Поэтому на первом слайде оставлены тема, авторские данные, актуальный экран каталога и короткая логика проекта.",
        "Поэтому на первом слайде оставлены только тема, авторские данные и актуальный экран каталога; подробная логика проекта раскрывается на следующих слайдах.",
    )
    text = text.replace(
        "актуальной светлой презентации на 12 слайдов",
        "актуальной финальной презентации на 12 слайдов",
    )
    md_path.write_text(text, encoding="utf-8")

    for docx_path in sorted(DOCS.glob("Речь_и_вопросы_к_защите_DOLG_*.docx"), key=lambda p: p.stat().st_mtime, reverse=True):
        if docx_path.name.startswith("~$"):
            continue
        doc = Document(str(docx_path))
        changed = False
        for paragraph in doc.paragraphs:
            if "актуальный экран каталога и короткая логика проекта" in paragraph.text:
                paragraph.text = paragraph.text.replace(
                    "актуальный экран каталога и короткая логика проекта",
                    "актуальный экран каталога; подробная логика проекта раскрывается на следующих слайдах",
                )
                changed = True
            if "актуальной светлой презентации на 12 слайдов" in paragraph.text:
                paragraph.text = paragraph.text.replace(
                    "актуальной светлой презентации на 12 слайдов",
                    "актуальной финальной презентации на 12 слайдов",
                )
                changed = True
        if changed:
            try:
                doc.save(str(docx_path))
            except PermissionError:
                fallback = docx_path.with_name(f"{docx_path.stem}_final_clean_20260519{docx_path.suffix}")
                doc.save(str(fallback))
        break


class FinalDeck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
        self.bg = RGBColor(250, 252, 255)
        self.white = RGBColor(255, 255, 255)
        self.navy = RGBColor(16, 33, 58)
        self.text = RGBColor(43, 55, 76)
        self.muted = RGBColor(93, 106, 126)
        self.line = RGBColor(209, 219, 231)
        self.cyan = RGBColor(0, 145, 181)
        self.blue = RGBColor(37, 99, 235)
        self.green = RGBColor(20, 140, 72)
        self.orange = RGBColor(217, 92, 17)
        self.purple = RGBColor(117, 67, 205)
        self.red = RGBColor(205, 50, 50)

    @staticmethod
    def e(value: float):
        return Inches(value)

    def slide(self, number: int):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.bg
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.e(16), self.e(0.13))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.cyan
        bar.line.fill.background()
        return slide

    def textbox(
        self,
        slide,
        text: str,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        size: float = 16,
        bold: bool = False,
        color: RGBColor | None = None,
        align: str = "left",
        font: str = "Arial",
    ):
        box = slide.shapes.add_textbox(self.e(x), self.e(y), self.e(w), self.e(h))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = self.e(0.03)
        frame.margin_right = self.e(0.03)
        frame.margin_top = self.e(0.02)
        frame.margin_bottom = self.e(0.02)
        frame.vertical_anchor = MSO_ANCHOR.TOP
        for idx, part in enumerate(text.split("\n")):
            paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            paragraph.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
            run = paragraph.add_run()
            run.text = part
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color or self.text
        return box

    def header(self, slide, title: str, subtitle: str, number: int) -> None:
        self.textbox(slide, title, 0.72, 0.42, 10.5, 0.45, size=26, bold=True, color=self.navy)
        if subtitle:
            self.textbox(slide, subtitle, 0.74, 0.94, 12.2, 0.36, size=15, color=self.muted)
        self.textbox(slide, f"{number}/12", 14.72, 0.55, 0.75, 0.32, size=13, color=self.muted, align="right")

    def panel(self, slide, x: float, y: float, w: float, h: float, accent: RGBColor | None = None):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.e(x), self.e(y), self.e(w), self.e(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.white
        shape.line.color.rgb = self.line
        shape.line.width = Pt(1)
        if accent:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.e(x), self.e(y), self.e(w), self.e(0.08))
            line.fill.solid()
            line.fill.fore_color.rgb = accent
            line.line.fill.background()
        return shape

    def card(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        title: str,
        body: str,
        accent: RGBColor,
        *,
        title_size: float = 16,
        body_size: float = 14,
    ) -> None:
        self.panel(slide, x, y, w, h, accent)
        self.textbox(slide, title, x + 0.22, y + 0.2, w - 0.44, 0.35, size=title_size, bold=True, color=self.navy)
        self.textbox(slide, body, x + 0.22, y + 0.72, w - 0.44, h - 0.82, size=body_size, color=self.text)

    def bullet_card(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        title: str,
        bullets: list[str],
        accent: RGBColor,
        *,
        title_size: float = 16,
        body_size: float = 14,
    ) -> None:
        self.panel(slide, x, y, w, h, accent)
        self.textbox(slide, title, x + 0.22, y + 0.2, w - 0.44, 0.35, size=title_size, bold=True, color=self.navy)
        body = "\n".join(f"- {item}" for item in bullets)
        self.textbox(slide, body, x + 0.25, y + 0.72, w - 0.5, h - 0.82, size=body_size, color=self.text)

    def metric(self, slide, x: float, y: float, w: float, h: float, value: str, label: str, accent: RGBColor) -> None:
        self.panel(slide, x, y, w, h, accent)
        self.textbox(slide, value, x + 0.04, y + 0.17, w - 0.08, 0.38, size=22, bold=True, color=accent, align="center")
        self.textbox(slide, label, x + 0.08, y + 0.65, w - 0.16, 0.36, size=13, color=self.muted, align="center")

    def fact(self, slide, x: float, y: float, w: float, h: float, text: str, accent: RGBColor) -> None:
        self.panel(slide, x, y, w, h, accent)
        self.textbox(slide, "Факт", x + 0.18, y + 0.16, 0.85, 0.3, size=13, bold=True, color=accent)
        self.textbox(slide, text, x + 1.0, y + 0.16, w - 1.16, h - 0.24, size=13.2, color=self.text)

    def image(self, slide, path: Path, x: float, y: float, w: float, h: float, caption: str = "") -> None:
        self.panel(slide, x, y, w, h + (0.36 if caption else 0), None)
        with Image.open(path) as image:
            iw, ih = image.size
        scale = min((w - 0.22) / iw, (h - 0.22) / ih)
        pw = iw * scale
        ph = ih * scale
        slide.shapes.add_picture(str(path), self.e(x + (w - pw) / 2), self.e(y + (h - ph) / 2), width=self.e(pw), height=self.e(ph))
        if caption:
            self.textbox(slide, caption, x + 0.15, y + h + 0.08, w - 0.3, 0.25, size=13, color=self.muted, align="center")


def build() -> Path:
    deck = FinalDeck()
    catalog = ASSETS / "catalog_cards_crop.jpg"
    lab = ASSETS / "engineering_lab_crop.jpg"
    cad = ASSETS / "cad_editor_crop.jpg"
    sim = ASSETS / "simulation_ac_crop.jpg"

    s = deck.slide(1)
    deck.textbox(s, "1/12", 14.72, 0.55, 0.75, 0.32, size=13, color=deck.muted, align="right")
    deck.textbox(s, "Дипломная работа", 0.75, 0.74, 6.0, 0.56, size=34, bold=True, color=deck.navy)
    deck.textbox(s, TITLE, 0.78, 1.62, 6.65, 1.72, size=19, bold=True, color=deck.text)
    deck.card(s, 0.78, 5.35, 6.15, 1.42, "Автор", f"{AUTHOR}\nНаучный руководитель: {SUPERVISOR}\nГод защиты: 2026", deck.cyan, title_size=15.5, body_size=13.8)
    deck.image(s, catalog, 8.08, 0.88, 7.02, 6.62, "Актуальный каталог и карточки компонентов")

    s = deck.slide(2)
    deck.header(s, "Проблема и цель", "Что именно решает DOLG", 2)
    deck.card(s, 0.8, 1.55, 4.55, 2.35, "Проблема", "Каталог, datasheet, CAD, SPICE, BOM и заказ часто существуют отдельно. Номиналы, модели и результаты приходится переносить вручную.", deck.blue, body_size=14.6)
    deck.card(s, 5.75, 1.55, 4.55, 2.35, "Цель", "Разработать веб-приложение для подбора и покупки компонентов со встроенными инструментами проектирования и симуляции схем.", deck.orange, body_size=14.6)
    deck.card(s, 10.7, 1.55, 4.55, 2.35, "Задачи", "Реализовать каталог, учетные записи, корзину и заказы, редактор схем, SPICE-симуляцию, лабораторию, обучение и проверки.", deck.green, body_size=14.6)
    for i, (value, label, color) in enumerate([("89", "товаров", deck.cyan), ("43", "РЭБ-компонента", deck.blue), ("5", "расчетов", deck.green), ("29", "заданий", deck.orange)]):
        deck.metric(s, 2.05 + i * 3.1, 5.25, 2.1, 1.22, value, label, color)

    s = deck.slide(3)
    deck.header(s, "Анализ аналогов", "Сильные стороны существующих решений и ниша проекта", 3)
    deck.card(s, 0.88, 1.55, 3.35, 2.0, "Маркетплейсы", "Сильны в поиске, цене и заказе, но не проверяют схему и инженерные расчеты.", deck.blue)
    deck.card(s, 4.48, 1.55, 3.35, 2.0, "Симуляторы", "Позволяют быстро экспериментировать, но обычно слабо связаны с реальными товарами.", deck.green)
    deck.card(s, 8.08, 1.55, 3.35, 2.0, "CAD/EDA", "KiCad и Altium дают полный workflow, но имеют высокий порог входа.", deck.orange)
    deck.card(s, 11.68, 1.55, 3.35, 2.0, "AI-CAD", "Дают подсказки, но требуют объяснимого источника инженерного вывода.", deck.purple)
    deck.card(s, 1.1, 4.8, 13.65, 1.25, "Позиция DOLG", "Учебно-инженерный веб-сервис: каталог, схема, расчет, измерение, обучение и заказ находятся в одном контуре.", deck.cyan, title_size=17, body_size=15.0)
    deck.fact(s, 1.1, 6.35, 13.65, 0.72, "Отличие проекта не в попытке заменить KiCad или Altium, а в связке: реальный компонент -> схема -> расчет -> проверка -> учебное объяснение.", deck.green)

    s = deck.slide(4)
    deck.header(s, "Целевая аудитория", "Для кого проект полезен", 4)
    deck.card(s, 0.9, 1.6, 4.35, 2.35, "Студенты", "Осваивают схемотехнику, решают практические задания и получают обратную связь по ошибкам.", deck.blue, body_size=15)
    deck.card(s, 5.82, 1.6, 4.35, 2.35, "Радиолюбители", "Подбирают компоненты, проверяют простые узлы и формируют набор деталей для сборки.", deck.green, body_size=15)
    deck.card(s, 10.74, 1.6, 4.35, 2.35, "Инженеры", "Сравнивают номиналы, оценивают запас, повторяемость решений и риски проекта.", deck.orange, body_size=15)
    deck.bullet_card(s, 2.05, 5.2, 11.9, 1.65, "Что объединяет аудитории", ["нужны реальные компоненты, а не абстрактные элементы", "важны проверяемые расчеты и понятные ошибки", "полезна связь между схемой, обучением и BOM"], deck.cyan, title_size=16, body_size=14.2)

    s = deck.slide(5)
    deck.header(s, "Архитектура системы", "Формулы и проверки вынесены в общий service-layer", 5)
    layers = [
        ("Интерфейс", "каталог, карточки товаров, редактор схем, симулятор, лаборатория, обучение"),
        ("Django apps", "shop, accounts, orders, knowledge, Dolg_APP, административный контур"),
        ("Service-layer", "engineering_lab, project_review, learning_grader, cad_import, rule_ai"),
        ("Expert-first core", "jsonschema, rule-engine, Pint, Lark, Z3, scikit-fuzzy"),
        ("Данные", "товары, РЭБ-компоненты, схемы, уроки, попытки, review snapshots"),
    ]
    colors = [deck.cyan, deck.blue, deck.green, deck.orange, deck.purple]
    for i, ((title, body), color) in enumerate(zip(layers, colors)):
        deck.card(s, 0.9, 1.45 + i * 1.18, 14.2, 0.86, title, body, color, title_size=15.5, body_size=13.6)
    deck.fact(s, 0.9, 7.32, 14.2, 0.62, "Новые библиотеки не тянутся в шаблоны: они подключены через сервисы, поэтому один расчет используется в лаборатории, обучении, review и AI-помощнике.", deck.cyan)

    s = deck.slide(6)
    deck.header(s, "Каталог и наполнение", "Товары связаны с документацией, статьями и инженерными сценариями", 6)
    deck.image(s, catalog, 0.78, 1.48, 8.9, 5.55, "Карточки товаров и актуальные изображения")
    deck.metric(s, 10.25, 1.65, 1.9, 1.15, "89", "товаров", deck.cyan)
    deck.metric(s, 12.65, 1.65, 1.9, 1.15, "43", "РЭБ", deck.blue)
    deck.metric(s, 10.25, 3.18, 1.9, 1.15, "21", "статья", deck.green)
    deck.metric(s, 12.65, 3.18, 1.9, 1.15, "50", "материалов", deck.orange)
    deck.bullet_card(s, 10.15, 5.0, 4.65, 1.55, "Медиа-политика", ["нелепые фото заменяются", "Wikimedia больше не основной источник", "приоритет - производитель и техпорталы"], deck.purple, body_size=13.5)
    deck.fact(s, 10.15, 6.9, 4.65, 0.72, "Карточка товара стала входной точкой для datasheet, статьи, схемы, BOM и заказа.", deck.cyan)

    s = deck.slide(7)
    deck.header(s, "CAD-редактор", "Схема хранится как данные для проверки и симуляции", 7)
    deck.image(s, cad, 0.78, 1.45, 10.0, 5.75, "Рабочая область редактора схем")
    deck.bullet_card(s, 11.15, 1.55, 3.95, 1.65, "Проверки", ["GND", "источник", "соединения", "номиналы"], deck.blue)
    deck.bullet_card(s, 11.15, 3.62, 3.95, 1.65, "Выходные данные", ["scheme_data", "BOM", "netlist", "review input"], deck.green)
    deck.bullet_card(s, 11.15, 5.68, 3.95, 1.35, "Смысл", ["схема не картинка", "это модель проекта"], deck.orange)
    deck.fact(s, 0.78, 7.48, 14.32, 0.56, "Одна и та же схема используется для DRC/ERC, генерации netlist, учебной проверки и будущего CAD-импорта.", deck.blue)

    s = deck.slide(8)
    deck.header(s, "Симуляция", "DC, AC и TRAN анализы дают измеримые результаты", 8)
    deck.image(s, sim, 0.78, 1.45, 9.65, 5.72, "Графики и результаты симуляции")
    deck.card(s, 10.8, 1.55, 4.15, 1.45, "Режимы", "DC / OP, AC, TRAN, Bode plot, FFT spectrum, Monte Carlo tolerance.", deck.blue, body_size=13.8)
    deck.card(s, 10.8, 3.3, 4.15, 1.45, "Измерения", "Напряжение узла, ток ветви, RMS, частота, duty cycle, мощность.", deck.green, body_size=13.8)
    deck.card(s, 10.8, 5.05, 4.15, 1.45, "Оценка", "Норма, риск, перегрев, нужен запас. Результат объясняется.", deck.orange, body_size=13.8)
    deck.fact(s, 0.78, 7.48, 14.32, 0.56, "Симуляция используется не как отдельный график, а как проверка гипотезы: расчетное значение сравнивается с измеренным.", deck.green)

    s = deck.slide(9)
    deck.header(s, "Инженерная лаборатория", "Расчеты сразу дают инженерную интерпретацию", 9)
    deck.image(s, lab, 0.78, 1.45, 8.8, 5.78, "Расчет и оценка результата")
    deck.bullet_card(s, 10.05, 1.55, 4.85, 2.0, "Расчеты", ["транзисторный ключ", "NE555", "стабилизатор", "тепловой запас", "RC-антидребезг"], deck.green, body_size=14.2)
    deck.bullet_card(s, 10.05, 4.02, 4.85, 1.85, "Использование", ["лаборатория", "обучающие задания", "Engineering Review"], deck.cyan, body_size=14.2)
    deck.fact(s, 10.05, 6.35, 4.85, 0.82, "Лаборатория возвращает не только число, но и инженерную оценку: норма, риск, перегрев или нужен запас.", deck.orange)

    s = deck.slide(10)
    deck.header(s, "Обучение и Engineering Review", "Практикум связан с реальными проверками проекта", 10)
    deck.card(s, 0.88, 1.5, 4.25, 2.05, "Практикум", "4 маршрута, 13 уроков, 29 заданий: численные ответы, сборка схемы, измерение результата.", deck.blue, body_size=14.3)
    deck.card(s, 5.55, 1.5, 4.25, 2.05, "Review", "Design Health Score, DRC/ERC, BOM-риски, наличие GND, источника, моделей и запаса.", deck.green, body_size=14.3)
    deck.card(s, 10.22, 1.5, 4.25, 2.05, "AI-помощник", "Самописный rule_ai отвечает по фактам проекта, правилам, расчетам, статьям и заданиям.", deck.orange, body_size=14.3)
    deck.bullet_card(s, 1.2, 4.95, 13.25, 1.65, "Expert-first подход", ["сначала объяснимые правила и факты", "затем constraint/optimization", "нейронная подсказка только поверх проверяемого baseline"], deck.cyan, title_size=16.5, body_size=14.5)
    deck.fact(s, 1.2, 6.95, 13.25, 0.62, "Учебное задание может строиться из реальной ошибки проекта: нет GND, неверный номинал, перегрев или расхождение расчет/симуляция.", deck.purple)

    s = deck.slide(11)
    deck.header(s, "Проверка и развитие", "Что подтверждено и куда проект развивается дальше", 11)
    checks = [("manage.py check", "0 issues"), ("demo-ready", "OK"), ("data-integrity", "OK"), ("targeted tests", "expert / learning / search")]
    for i, (title, body) in enumerate(checks):
        deck.card(s, 0.9 + i * 3.65, 1.55, 3.15, 1.45, title, body, [deck.green, deck.cyan, deck.blue, deck.orange][i], title_size=15, body_size=14.2)
    deck.bullet_card(s, 0.95, 4.05, 6.8, 2.05, "Ближайший этап", ["расширение Engineering Review", "сохраненные измерения", "CAD-import subset", "комментарии пользователей"], deck.green, title_size=16.5, body_size=14.4)
    deck.bullet_card(s, 8.15, 4.05, 6.8, 2.05, "Поздний этап", ["OR-Tools для BOM-оптимизации", "RDFLib для онтологии компонентов", "PyTorch/GOLEM для deep hints"], deck.purple, title_size=16.5, body_size=14.4)

    s = deck.slide(12)
    deck.textbox(s, "12/12", 14.72, 0.55, 0.75, 0.32, size=13, color=deck.muted, align="right")
    deck.textbox(s, "Спасибо за внимание!", 0.9, 1.05, 8.1, 0.75, size=36, bold=True, color=deck.navy)
    deck.textbox(s, "DOLG объединяет каталог, CAD/SIM, инженерную лабораторию, обучение и экспертную проверку проекта.", 0.94, 2.05, 7.65, 0.9, size=18, color=deck.text)
    deck.card(s, 1.0, 4.25, 5.85, 1.55, "Контакты", f"{AUTHOR}\nEmail: buryako@internet.com\nТелефон: +7 (903) 439-44-87", deck.cyan, title_size=15.5, body_size=13.8)
    deck.card(s, 7.55, 4.25, 6.8, 1.55, "Материалы", "Диплом, речь и презентация актуализированы. Проверки проекта проходят.", deck.green, title_size=15.5, body_size=13.8)
    deck.textbox(s, "Вопросы?", 5.3, 7.0, 5.4, 0.6, size=28, bold=True, color=deck.cyan, align="center")

    docs_pptx = DOCS / "Презентация_DOLG_финальная_20260513_v5.pptx"
    actual_pptx = DOCS / "Презентация_DOLG_актуальная_20260519.pptx"
    final_pptx = DOCS / "Презентация_DOLG_финальная_читабельная_20260519.pptx"
    for path in [docs_pptx, actual_pptx, final_pptx]:
        deck.prs.save(str(path))

    download_target = DOWNLOADS / "Razrabotka-veb-prilozheniya-dlya-prodazhi-radio-i-elektronnyh-komponentov-so-vstroennymi-instrumenta_updated_20260518.pptx"
    download_final = DOWNLOADS / "Razrabotka-veb-prilozheniya-dlya-prodazhi-radio-i-elektronnyh-komponentov-so-vstroennymi-instrumenta_final_clean_20260519.pptx"
    if download_target.exists():
        copy_safe(docs_pptx, download_target)
    copy_safe(docs_pptx, download_final)

    Presentation(str(docs_pptx))
    Presentation(str(actual_pptx))
    Presentation(str(final_pptx))
    return final_pptx


def main() -> None:
    backup_dir = DOCS / "presentation_backups" / "20260519_final_clean_before_script"
    backup_dir.mkdir(parents=True, exist_ok=True)
    for pattern in ["Презентация_DOLG_*.pptx", "Речь_и_вопросы_к_защите_DOLG_*.md", "Речь_и_вопросы_к_защите_DOLG_*.docx"]:
        for path in DOCS.glob(pattern):
            if not path.name.startswith("~$"):
                backup(path, backup_dir)
    download_pptx = DOWNLOADS / "Razrabotka-veb-prilozheniya-dlya-prodazhi-radio-i-elektronnyh-komponentov-so-vstroennymi-instrumenta_updated_20260518.pptx"
    backup(download_pptx, backup_dir)

    update_speech_note()
    out = build()
    print(f"Built final clean presentation: {out}")


if __name__ == "__main__":
    main()
